from __future__ import annotations

import argparse
import hashlib
import os
import sys
import time
import re
import warnings
from datetime import date, datetime, timedelta
from bisect import bisect_right
from dataclasses import dataclass
from pathlib import Path
from collections import defaultdict
from typing import Callable, Dict, Iterable, List, Optional, Set, Tuple

COMMON_MANUAL_RULES: list = []  # 例外防止のため空リストで初期化（必要ならルールを追加）
# --- 簡易プロファイリング機能 ---
import functools
PROFILE_ENABLED = False  # Trueで有効化
_profile_results = []

def profile_section(label):
    def decorator(func):
        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            if not PROFILE_ENABLED:
                return func(*args, **kwargs)
            start = time.perf_counter()
            result = func(*args, **kwargs)
            elapsed = time.perf_counter() - start
            _profile_results.append((label, elapsed))
            return result
        return wrapper
    return decorator

def print_profile_results():
    if not PROFILE_ENABLED:
        return
    print("\n[PROFILE] major section timings:")
    for label, elapsed in _profile_results:
        print(f"  {label}: {elapsed:.3f} sec")


try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.utils.cell import coordinate_to_tuple, get_column_letter, range_boundaries
except Exception:
    Workbook = None
    XLImage = None
    load_workbook = None
    coordinate_to_tuple = None
    get_column_letter = None
    range_boundaries = None

# Pillow依存部のimport（Excelフォールバック画像化用）
try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:
    Image = None
    ImageDraw = None
    ImageFont = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import fitz
except Exception:
    fitz = None

try:
    from win32com.client import gencache
except Exception:
    gencache = None

MM_PER_INCH = 25.4
EMUS_PER_INCH = 914400
EMU_PER_MM = EMUS_PER_INCH / MM_PER_INCH
COMMON_MARGIN_RULE_TEXT = "余白（A4縦:上20/下20/左30/右20mm以上）"

EXCLUDED_CHECK_IDS: Set[str] = set()
EXCLUDED_CHECK_ITEMS_NORMALIZED: Set[str] = set()


def normalize_check_item_key(check_item: object) -> str:
    text = str(check_item or "").strip().lower()
    text = re.sub(r"\s+", "", text)
    for prefix in ("共通:", "word:", "excel:", "pdf:", "visio:"):
        if text.startswith(prefix):
            return text[len(prefix):]
    return text



def is_excluded_check(check_id: object = "", check_item: object = "") -> bool:
    # C3（余白チェック）は完全除外
    if str(check_id).strip().upper() == "C3":
        return True
    return False


DEFAULT_COVER_KEYWORDS: List[str] = [
    "進捗状況報告",
    "進捗報告",
    "月間会議計画",
    "マスタースケジュール",
    "EVM",
    "リスク管理表",
    "別紙",
]


_JP_HOLIDAY_CACHE: Dict[int, Dict[date, str]] = {}
_WORD_PAGE_MAP_CACHE: Dict[str, Dict[str, List[int]]] = {}
SUGGESTED_ACTION_SETTINGS_SHEET = "suggested_action_settings"


@dataclass
class SuggestedActionSetting:
    enabled: bool
    file_type: str
    check_id: str
    status: str
    suggested_action: str
    check_item: str = ""
    excel_status: str = ""

@dataclass
class CheckResult:
    file_path: str
    file_type: str
    check_id: str
    check_item: str
    status: str
    detail: str
    suggested_action: str

@dataclass
class VisualPage:
    file_path: str
    page_no: int
    image_path: str
    sheet_name: str = ""

COMMON_CHECK_ITEMS: List[Tuple[str, str]] = [
    ("C1", "プロパティ情報削除"),
    ("C2", "表紙が規定のもの"),
    ("C3", COMMON_MARGIN_RULE_TEXT),
    ("C4", "ページ番号"),
    ("C5", "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）"),
]

TYPE_SPECIFIC_CHECK_ITEMS = {
    "Excel": [
        ("E1", "印刷範囲外記載チェック"),
        ("E2", "不要シートチェック"),
        ("E3", "印刷プレビュー（見切れ/罫線欠け）"),
        ("E4", "見え消し（取り消し線/訂正線）残存"),
        ("E5", "コメント（吹き出し）残存"),
        ("E6", "参照エラー残存"),
        ("E7", "ページ設定のタイトル行"),
    ],
    "Word": [
        ("W1", "見え消し（変更履歴）残存"),
        ("W2", "マーカ残存"),
        ("W3", "不要な文字色"),
        ("W4", "コメント（吹き出し）残存"),
        ("W5", "参照エラー残存"),
    ],
    "PDF": [
        ("P1", "見え消し（注釈ベース）残存"),
        ("P2", "コメント（吹き出し）残存"),
        ("P3", "参照エラー残存"),
    ],
    "PPT": [],
    "VISIO": [],
    "LegacyOffice": [("L1", "旧形式ファイル")],
}


DOMAIN_RULES: Dict[str, List[str]] = {}


DOMAIN_KEYWORDS: Dict[str, List[str]] = {}


def detect_domains(file_path: Path) -> List[str]:
    name = file_path.stem.lower()
    domains: List[str] = []
    for domain, words in DOMAIN_KEYWORDS.items():
        if any(w.lower() in name for w in words):
            domains.append(domain)
    return domains





def inches_to_mm(value: Optional[float]) -> Optional[float]:
    if value is None:
        return None
    try:
        return float(value) * MM_PER_INCH
    except Exception:
        return None


def _margin_orientation_label(is_landscape: bool) -> str:
    return "横" if is_landscape else "縦"


def _judge_c3_margin(top: Optional[float], bottom: Optional[float], left: Optional[float], right: Optional[float], is_landscape: bool) -> Tuple[bool, str]:
    """
    C3判定:
      縦設定: 左30mm以上、右・上・下20mm以上
      横設定: 上30mm以上、左・右・下20mm以上
    """
    if None in (top, bottom, left, right):
        return False, "余白取得不可"
    if is_landscape:
        ok = top >= 30 and left >= 20 and right >= 20 and bottom >= 20
        basis = "横設定: 上30mm以上、左・右・下20mm以上"
    else:
        ok = left >= 30 and top >= 20 and right >= 20 and bottom >= 20
        basis = "縦設定: 左30mm以上、右・上・下20mm以上"
    label = _margin_orientation_label(is_landscape)
    judgement = f"{label}設定として判定〇" if ok else f"{label}設定として判定×"
    return ok, f"上{top:.1f}/下{bottom:.1f}/左{left:.1f}/右{right:.1f}mm / {basis} / {judgement}"


def _excel_is_landscape(ws) -> bool:
    orientation = str(getattr(getattr(ws, "page_setup", None), "orientation", "") or "").lower()
    if "landscape" in orientation or orientation in {"横", "yoko"}:
        return True
    if "portrait" in orientation or orientation in {"縦", "tate"}:
        return False
    try:
        paper_width = getattr(ws.page_setup, "paperWidth", None)
        paper_height = getattr(ws.page_setup, "paperHeight", None)
        if paper_width and paper_height:
            return float(paper_width) > float(paper_height)
    except Exception:
        pass
    return False


def _word_is_landscape(section) -> bool:
    try:
        width = mm_from_emu(getattr(section, "page_width", None))
        height = mm_from_emu(getattr(section, "page_height", None))
        if width is not None and height is not None and width > height:
            return True
    except Exception:
        pass
    orientation = str(getattr(section, "orientation", "") or "").lower()
    return "landscape" in orientation or orientation in {"1", "wdorientlandscape", "横"}


def _runtime_base_dirs() -> List[Path]:
    """Windows配布向けに、実行ファイル周辺とPyInstaller展開先を探索する。"""
    dirs: List[Path] = []
    try:
        if getattr(sys, "frozen", False):
            dirs.append(Path(sys.executable).resolve().parent)
        dirs.append(Path(__file__).resolve().parent)
        meipass = getattr(sys, "_MEIPASS", None)
        if meipass:
            dirs.append(Path(meipass).resolve())
    except Exception:
        pass

    unique: List[Path] = []
    seen: Set[str] = set()
    for d in dirs:
        key = str(d).lower()
        if key not in seen:
            unique.append(d)
            seen.add(key)
    return unique


def _add_aspose_plugin_paths() -> List[str]:
    """
    オフラインWindows配布向けのAspose.Diagramプラグイン探索。
    pip自動導入は行わず、実行ファイル配下の plugins/vendor/lib に配置された
    aspose パッケージまたは site-packages 相当フォルダを sys.path に追加する。
    """
    candidates: List[Path] = []
    env_dir = os.environ.get("ASPOSE_DIAGRAM_PLUGIN_DIR", "").strip()
    if env_dir:
        candidates.append(Path(env_dir))

    for base_dir in _runtime_base_dirs():
        candidates.extend([
            base_dir / "plugins" / "aspose_diagram",
            base_dir / "plugins" / "aspose",
            base_dir / "plugins",
            base_dir / "vendor" / "aspose_diagram",
            base_dir / "vendor" / "aspose",
            base_dir / "vendor",
            base_dir / "lib" / "aspose_diagram",
            base_dir / "lib",
            base_dir / "_internal" / "plugins" / "aspose_diagram",
            base_dir / "_internal" / "plugins",
        ])

    added: List[str] = []
    for path in candidates:
        try:
            if path.exists() and path.is_dir():
                path_str = str(path.resolve())
                if path_str not in sys.path:
                    sys.path.insert(0, path_str)
                    added.append(path_str)
        except Exception:
            continue
    return added


def ensure_aspose_diagram_available(auto_install: bool = False) -> Tuple[bool, str]:
    """
    Aspose.Diagramを利用可能にする。
    オフライン・Python未導入PCでの配布を想定し、pip installは行わない。
    配布時は exe と同階層の plugins\aspose_diagram などに aspose パッケージを配置する。
    """
    try:
        import aspose.diagram  # type: ignore
        return True, "Aspose.Diagram 利用可能"
    except Exception as first_exc:
        added = _add_aspose_plugin_paths()
        try:
            import importlib
            importlib.invalidate_caches()
            import aspose.diagram  # type: ignore
            return True, "Aspose.Diagram プラグインを読み込みました。"
        except Exception as plugin_exc:
            hint = (
                "Aspose.Diagram プラグイン未配置または読み込み不可。"
                "exeと同じフォルダ配下の plugins\\aspose_diagram に aspose パッケージ一式を配置してください。"
            )
            if added:
                hint += f" 探索済み={'; '.join(added[:3])}"
            return False, f"{hint} import_error={first_exc}; plugin_error={plugin_exc}"


def add_result(
    results: List[CheckResult],
    file_path: Path,
    file_type: str,
    check_id: str,
    check_item: str,
    status: str,
    detail: str,
    suggested_action: str,
) -> None:
    if is_excluded_check(check_id, check_item):
        return
    results.append(
        CheckResult(
            file_path=str(file_path),
            file_type=file_type,
            check_id=check_id,
            check_item=check_item,
            status=status,
            detail=detail,
            suggested_action=suggested_action,
        )
    )


def _normalize_setting_status(value: object) -> str:
    text = str(value or "").strip().upper()
    if text in {"×", "X", "NG"}:
        return "CROSS"
    if text in {"〇", "○", "O", "OK", "PASS"}:
        return "PASS"
    if text in {"要チェック", "MANUAL", "N/A", "WARN"}:
        return "CHECK"
    if text == "FAIL":
        return "FAIL"
    if text == "ERROR":
        return "ERROR"
    return text


def _parse_bool_like(value: object, default: bool = True) -> bool:
    if value is None:
        return default
    text = str(value).strip().lower()
    if not text:
        return default
    if text in {"1", "true", "yes", "y", "on", "enabled", "有効"}:
        return True
    if text in {"0", "false", "no", "n", "off", "disabled", "無効"}:
        return False
    return default


def load_suggested_action_settings(
    config_xlsx: Optional[Path],
) -> Tuple[Dict[Tuple[str, str, str], str], List[SuggestedActionSetting], Optional[str]]:
    if config_xlsx is None:
        return {}, [], None
    if load_workbook is None:
        return {}, [], f"対応推奨設定を読み込めません。openpyxl が未インストールです: {config_xlsx}"
    if not config_xlsx.exists():
        return {}, [], f"対応推奨設定Excelが見つかりません: {config_xlsx}"

    try:
        wb = load_workbook(config_xlsx, data_only=True, read_only=True)
    except Exception as exc:
        return {}, [], f"対応推奨設定Excelの読み込みに失敗しました: {config_xlsx} ({exc})"

    try:
        if SUGGESTED_ACTION_SETTINGS_SHEET not in wb.sheetnames:
            return {}, [], None

        ws = wb[SUGGESTED_ACTION_SETTINGS_SHEET]
        header_cells = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), None)
        if not header_cells:
            return {}, [], None

        headers = {str(cell).strip(): idx for idx, cell in enumerate(header_cells) if cell is not None and str(cell).strip()}
        required = {"file_type", "check_id", "status", "suggested_action"}
        if not required.issubset(headers):
            missing = ", ".join(sorted(required - set(headers)))
            return {}, [], f"対応推奨設定シートの列が不足しています: {missing}"

        overrides: Dict[Tuple[str, str, str], str] = {}
        settings: List[SuggestedActionSetting] = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            file_type = str(row[headers["file_type"]] or "").strip()
            check_id = str(row[headers["check_id"]] or "").strip()
            status = _normalize_setting_status(row[headers["status"]])
            suggested_action = str(row[headers["suggested_action"]] or "").strip()
            enabled = _parse_bool_like(row[headers["enabled"]], default=True) if "enabled" in headers else True
            check_item = str(row[headers["check_item"]] or "").strip() if "check_item" in headers else ""
            excel_status = str(row[headers["excel_status"]] or "").strip() if "excel_status" in headers else ""

            if not check_id:
                continue

            setting = SuggestedActionSetting(
                enabled=enabled,
                file_type=file_type,
                check_id=check_id,
                status=status,
                suggested_action=suggested_action,
                check_item=check_item,
                excel_status=excel_status,
            )
            settings.append(setting)
            if enabled and suggested_action and status:
                overrides[(file_type.upper(), check_id.upper(), status)] = suggested_action

        return overrides, settings, None
    finally:
        try:
            wb.close()
        except Exception:
            pass


def _lookup_suggested_action_override(
    result: CheckResult,
    overrides: Dict[Tuple[str, str, str], str],
) -> Optional[str]:
    file_type = (result.file_type or "").upper()
    check_id = (result.check_id or "").upper()
    status_candidates: List[str] = []
    for candidate in (
        _normalize_setting_status(result.status),
        _normalize_setting_status(display_status(result.status)),
    ):
        if candidate and candidate not in status_candidates:
            status_candidates.append(candidate)

    candidates: List[Tuple[str, str, str]] = []
    for status in status_candidates:
        candidates.extend(
            [
                (file_type, check_id, status),
                ("", check_id, status),
            ]
        )
    for key in candidates:
        action = overrides.get(key)
        if action:
            return action
    return None


def apply_suggested_action_overrides(
    results: Iterable[CheckResult],
    overrides: Dict[Tuple[str, str, str], str],
) -> None:
    if not overrides:
        return
    for result in results:
        override = _lookup_suggested_action_override(result, overrides)
        if override:
            result.suggested_action = override


def build_suggested_action_settings(
    results: Iterable[CheckResult],
    existing_settings: Optional[Iterable[SuggestedActionSetting]] = None,
) -> List[SuggestedActionSetting]:
    merged: List[SuggestedActionSetting] = []
    seen: Set[Tuple[str, str, str]] = set()

    for setting in existing_settings or []:
        key = (setting.file_type.upper(), setting.check_id.upper(), setting.status)
        if key in seen:
            continue
        merged.append(setting)
        seen.add(key)

    for result in results:
        if display_status(result.status) != "×":
            continue
        key = ((result.file_type or "").upper(), (result.check_id or "").upper(), _normalize_setting_status(result.status))
        if key in seen:
            continue
        merged.append(
            SuggestedActionSetting(
                enabled=True,
                file_type=result.file_type,
                check_id=result.check_id,
                status=_normalize_setting_status(result.status),
                suggested_action=result.suggested_action,
                check_item=result.check_item,
                excel_status=display_status(result.status),
            )
        )
        seen.add(key)

    merged.sort(key=lambda s: ((s.file_type or "").upper(), (s.check_id or "").upper(), s.status, s.check_item or ""))
    return merged


def _normalize_property_text(value: object) -> str:
    if value is None:
        return "(空)"
    text = str(value).strip()
    return text if text else "(空)"


def _build_property_detail(fields: List[Tuple[str, object]]) -> Tuple[bool, str]:
    parts: List[str] = []
    has_any = False
    for label, value in fields:
        text = _normalize_property_text(value)
        if text != "(空)":
            has_any = True
        parts.append(f"{label}={text}")
    head = "プロパティ情報が残っています。" if has_any else "主要プロパティは空です。"
    return has_any, " / ".join([head] + parts)




def ensure_expected_checks(

    results: List[CheckResult],

    file_path: Path,

    file_type: str,

) -> None:

    target = str(file_path)

    present_ids = {r.check_id for r in results if r.file_path == target}

    expected = TYPE_SPECIFIC_CHECK_ITEMS.get(file_type, []) + COMMON_CHECK_ITEMS

    for check_id, check_item in expected:

        if check_id in present_ids:

            continue

        add_result(

            results,

            file_path,

            file_type,

            check_id,

            check_item,

            "N/A",

            "このファイルでは自動判定未対応、または判定対象外です。",

            "必要に応じて手動確認してください。",

        )



def slugify_for_path(path: Path) -> str:
    base = re.sub(r"[^A-Za-z0-9_.-]+", "_", path.stem)
    parent = re.sub(r"[^A-Za-z0-9_.-]+", "_", path.parent.name or "root")
    short_parent = parent[:16] or "root"
    short_base = base[:24] or "file"
    hash8 = hashlib.sha1(str(path).encode("utf-8", errors="ignore")).hexdigest()[:8]
    return f"{short_parent}_{short_base}_{hash8}"


@profile_section("convert_office_to_pdf")
def convert_office_to_pdf(input_path: Path, output_pdf: Path) -> Optional[str]:
    if gencache is None:
        return "pywin32 が未インストールのため Office→PDF 変換できません。"

    import pythoncom
    suffix = input_path.suffix.lower()
    if suffix in {".doc", ".docx"}:
        word = None
        doc = None
        pythoncom.CoInitialize()
        try:
            word = gencache.EnsureDispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(str(input_path), ReadOnly=True)
            doc.ExportAsFixedFormat(str(output_pdf), 17)
            return None
        except Exception as exc:
            return f"Word変換失敗: {exc}"
        finally:
            if doc is not None:
                try:
                    doc.Close(False)
                except Exception:
                    pass
            if word is not None:
                try:
                    word.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    if suffix in {".xls", ".xlsx", ".xlsm"}:
        excel = None
        wb = None
        pythoncom.CoInitialize()
        try:
            excel = gencache.EnsureDispatch("Excel.Application")
            excel.Visible = False
            excel.DisplayAlerts = False
            wb = excel.Workbooks.Open(str(input_path))
            wb.ExportAsFixedFormat(0, str(output_pdf))
            return None
        except Exception as exc:
            return f"Excel変換失敗: {exc}"
        finally:
            if wb is not None:
                try:
                    wb.Close(False)
                except Exception:
                    pass
            if excel is not None:
                try:
                    excel.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    return f"未対応拡張子です: {suffix}"


def convert_visio_to_pdf_via_aspose(input_path: Path, output_pdf: Path) -> Optional[str]:
    ok, msg = ensure_aspose_diagram_available(auto_install=False)
    if not ok:
        return f"Aspose.Diagram を利用できません(Python {sys.version_info.major}.{sys.version_info.minor}): {msg}"
    try:
        from aspose.diagram import Diagram, SaveFileFormat
    except Exception as exc:
        return f"Aspose.Diagram を利用できません(Python {sys.version_info.major}.{sys.version_info.minor}): {exc}"

    try:
        diagram = Diagram(str(input_path))
        output_pdf.parent.mkdir(parents=True, exist_ok=True)
        diagram.save(str(output_pdf), SaveFileFormat.PDF)
        return None
    except Exception as exc:
        return f"Aspose.Diagram 変換失敗: {exc}"


def convert_visio_to_pdf_via_com(input_path: Path, output_pdf: Path) -> Optional[str]:
    if gencache is None:
        return "pywin32 が未インストールのため Visio COM を利用できません。"

    import pythoncom

    visio = None
    docs = None
    pythoncom.CoInitialize()
    try:
        visio = gencache.EnsureDispatch("Visio.Application")
        visio.Visible = False
        docs = visio.Documents.Open(str(input_path))
        docs.ExportAsFixedFormat(1, str(output_pdf), 1, 0)
        return None
    except Exception as exc:
        return f"Visio COM 変換失敗: {exc}"
    finally:
        if docs is not None:
            try:
                docs.Close()
            except Exception:
                pass
        if visio is not None:
            try:
                visio.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def find_soffice_path() -> Optional[str]:
    candidates = [
        os.environ.get("SOFFICE_PATH", ""),
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def convert_visio_to_pdf_via_libreoffice(input_path: Path, output_pdf: Path) -> Optional[str]:
    import subprocess

    soffice_path = find_soffice_path()
    if soffice_path is None:
        return "LibreOffice が見つからないため変換できません。"

    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    try:
        result = subprocess.run(
            [
                soffice_path,
                "--headless",
                "--convert-to",
                "pdf:draw_pdf_Export",
                str(input_path),
                "--outdir",
                str(output_pdf.parent),
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
    except Exception as exc:
        return f"LibreOffice 変換失敗: {exc}"

    generated_pdf = output_pdf.parent / f"{input_path.stem}.pdf"
    if result.returncode != 0 and not generated_pdf.exists():
        detail = (result.stderr or result.stdout or "").strip()
        return f"LibreOffice 変換失敗: {detail or ('exit=' + str(result.returncode))}"

    if not generated_pdf.exists():
        detail = (result.stderr or result.stdout or "").strip()
        return f"LibreOffice で PDF が生成されませんでした。{(' detail=' + detail) if detail else ''}"

    if generated_pdf != output_pdf:
        try:
            generated_pdf.replace(output_pdf)
        except Exception as exc:
            return f"LibreOffice 生成PDF移動失敗: {exc}"
    return None


def render_visio_with_pdf_backend(
    file_path: Path,
    pdf_path: Path,
    image_dir: Path,
) -> Tuple[List[Path], Optional[str], Optional[str]]:
    if file_path.suffix.lower() == ".vsd":
        backends = [
            ("Visio COM", convert_visio_to_pdf_via_com),
            ("Aspose.Diagram", convert_visio_to_pdf_via_aspose),
            ("LibreOffice", convert_visio_to_pdf_via_libreoffice),
        ]
    else:
        backends = [
            ("Aspose.Diagram", convert_visio_to_pdf_via_aspose),
            ("Visio COM", convert_visio_to_pdf_via_com),
            ("LibreOffice", convert_visio_to_pdf_via_libreoffice),
        ]
    errors: List[str] = []

    for backend_name, converter in backends:
        err = converter(file_path, pdf_path)
        if err:
            errors.append(err)
            continue
        if not pdf_path.exists():
            errors.append(f"{backend_name} で PDF が生成されませんでした。")
            continue
        images, image_err = render_pdf_to_pngs(pdf_path, image_dir)
        if image_err:
            return [], image_err, backend_name
        return images, None, backend_name

    return [], " / ".join(errors), None


@profile_section("render_pdf_to_pngs")
def render_pdf_to_pngs(pdf_path: Path, image_dir: Path) -> Tuple[List[Path], Optional[str]]:
    if fitz is None:
        return [], "PyMuPDF が未インストールのため PDF→PNG 変換できません。"

    try:
        tools = getattr(fitz, "TOOLS", None)
        if tools is not None:
            if hasattr(tools, "mupdf_display_warnings"):
                tools.mupdf_display_warnings(False)
            if hasattr(tools, "mupdf_display_errors"):
                tools.mupdf_display_errors(False)
    except Exception:
        pass

    image_dir.mkdir(parents=True, exist_ok=True)
    try:
        doc = fitz.open(str(pdf_path))
    except Exception as exc:
        return [], f"PDF画像化失敗: {exc}"

    out_files: List[Path] = []
    try:
        for page_no, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            out = image_dir / f"page_{page_no:03d}.png"
            pix.save(str(out))
            out_files.append(out)
    except Exception as exc:
        return out_files, f"ページ画像化中に失敗: {exc}"
    finally:
        doc.close()

    return out_files, None


def sanitize_filename_for_path(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_.-]+", "_", name or "sheet")
    return cleaned[:40] or "sheet"


def is_excel_no_print_target_error(exc: Exception) -> bool:
    msg = str(exc)
    return ("印刷する対象がありません" in msg) or ("-2146827284" in msg)


def summarize_excel_sheet_export_error(exc: Exception) -> str:
    msg = str(exc)
    if "-2147024809" in msg or "無効な引数" in msg or "Invalid argument" in msg:
        return "COM無効引数"
    if is_excel_no_print_target_error(exc):
        return "印刷対象なし"
    return "COM例外"


def append_skipped_sheet(skipped_sheets: List[str], sheet_name: str, exc: Optional[Exception] = None) -> None:
    reason = "COM例外"
    if exc is not None:
        try:
            reason = summarize_excel_sheet_export_error(exc)
        except Exception:
            reason = "COM例外"
    try:
        skipped_sheets.append(f"{sheet_name}({reason})")
    except Exception:
        skipped_sheets.append(str(sheet_name))


def append_sheet(skipped_sheets: List[str], sheet_name: str, exc: Optional[Exception] = None) -> None:
    append_skipped_sheet(skipped_sheets, sheet_name, exc)


def render_excel_sheets_to_png_fallback(
    input_path: Path,
    image_dir: Path,
    max_rows: int = 120,
    max_cols: int = 18,
) -> Tuple[List[Tuple[str, Path]], Optional[str]]:
    if load_workbook is None:
        return [], "openpyxl が未インストールのため Excelフォールバック画像化できません。"
    if Image is None or ImageDraw is None:
        return [], "Pillow が未インストールのため Excelフォールバック画像化できません。"

    image_dir.mkdir(parents=True, exist_ok=True)
    rendered: List[Tuple[str, Path]] = []

    try:
        wb = load_workbook(input_path, data_only=True, read_only=False)
    except Exception as exc:
        return [], f"Excel読み込み失敗(フォールバック): {exc}"

    try:
        for idx, ws in enumerate(wb.worksheets, start=1):
            sheet_name = str(getattr(ws, "title", f"Sheet{idx}"))
            lines: List[str] = []
            added = 0

            for r in range(1, min(int(getattr(ws, "max_row", 1) or 1), max_rows) + 1):
                row_texts: List[str] = []
                for c in range(1, min(int(getattr(ws, "max_column", 1) or 1), max_cols) + 1):
                    value = ws.cell(row=r, column=c).value
                    if value is None:
                        row_texts.append("")
                        continue
                    text = str(value).replace("\r", " ").replace("\n", " ").strip()
                    row_texts.append(text[:40])
                if any(t for t in row_texts):
                    lines.append(" | ".join(row_texts).rstrip(" |"))
                    added += 1
                if added >= max_rows:
                    break

            if not lines:
                lines.append("(表示可能なセル値がありません)")

            font_w = 8
            line_h = 18
            padding = 12
            max_len = max(len(x) for x in lines)
            width = min(max(640, padding * 2 + max_len * font_w), 3800)
            height = min(max(220, padding * 2 + len(lines) * line_h), 7000)

            # --- 日本語フォント指定（優先: メイリオ, 次: MSゴシック, フォールバック: デフォルト） ---
            try:
                font = ImageFont.truetype("meiryo.ttc", 16)
            except Exception:
                try:
                    font = ImageFont.truetype("msgothic.ttc", 16)
                except Exception:
                    font = ImageFont.load_default()

            img = Image.new("RGB", (width, height), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            y = padding

            for line in lines:
                try:
                    draw.text((padding, y), line, fill=(0, 0, 0), font=font)
                except Exception:
                    # フォント描画でエラー時はデフォルトフォントで再描画
                    draw.text((padding, y), line, fill=(0, 0, 0))
                y += line_h

            out_path = image_dir / f"sheet_{idx:03d}_{sanitize_filename_for_path(sheet_name)}.png"
            img.save(out_path)
            rendered.append((sheet_name, out_path))
    finally:
        try:
            wb.close()
        except Exception:
            pass

    return rendered, None


@profile_section("convert_excel_to_sheet_pdfs")
def convert_excel_to_sheet_pdfs(input_path: Path, pdf_dir: Path, slug: str) -> Tuple[List[Tuple[str, Path]], List[str], Optional[str]]:
    if gencache is None:
        return [], [], "pywin32 が未インストールのため Excelシート別PDF変換できません。"

    excel = None
    wb = None
    out: List[Tuple[str, Path]] = []
    skipped_sheets: List[str] = []
    import pythoncom
    pythoncom.CoInitialize()
    try:
        excel = gencache.EnsureDispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        wb = excel.Workbooks.Open(str(input_path), ReadOnly=True)

        sheet_count = int(getattr(wb.Worksheets, "Count", 0) or 0)
        for idx in range(1, sheet_count + 1):
            sheet_name = f"Sheet{idx}"
            try:
                ws = wb.Worksheets(idx)
                sheet_name = str(getattr(ws, "Name", sheet_name))
            except Exception as exc:
                append_skipped_sheet(skipped_sheets, sheet_name, exc)
                continue
            safe_sheet = sanitize_filename_for_path(sheet_name)
            # TODO: VSDファイルの場合は将来的に .pdf → .vsd でExcelセル出力するよう拡張する（現状は.pdfのまま）
            pdf_path = pdf_dir / f"{slug}_s{idx:03d}_{safe_sheet}.pdf"
            if pdf_path.exists():
                try:
                    pdf_path.unlink()
                except Exception:
                    pass
            try:
                ws.ExportAsFixedFormat(0, str(pdf_path))
            except Exception:
                try:
                    ws.ExportAsFixedFormat(0, str(pdf_path), IgnorePrintAreas=True)
                except Exception as retry_exc:
                    append_skipped_sheet(skipped_sheets, sheet_name, retry_exc)
                    continue
            if pdf_path.exists():
                out.append((sheet_name, pdf_path))
            else:
                skipped_sheets.append(sheet_name)

        if not out:
            fallback_pdf = pdf_dir / f"{slug}_workbook.pdf"
            try:
                wb.ExportAsFixedFormat(0, str(fallback_pdf))
                if fallback_pdf.exists():
                    return [("(Workbook)", fallback_pdf)], skipped_sheets, None
            except Exception:
                try:
                    wb.ExportAsFixedFormat(0, str(fallback_pdf), IgnorePrintAreas=True)
                    if fallback_pdf.exists():
                        return [("(Workbook)", fallback_pdf)], skipped_sheets, None
                except Exception as retry_exc:
                    return [], skipped_sheets, f"Excelシート別PDF変換に失敗しました（出力0件）: {retry_exc}"
            return [], skipped_sheets, "Excelシート別PDF変換に失敗しました（出力0件）。"
        return out, skipped_sheets, None
    except Exception as exc:
        return out, skipped_sheets, f"Excelシート別変換失敗: {exc}"
    finally:
        if wb is not None:
            try:
                wb.Close(False)
            except Exception:
                pass
        if excel is not None:
            try:
                excel.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def main_with_profile(*args, **kwargs):
    try:
        result = main(*args, **kwargs)
    finally:
        print_profile_results()
    return result

@profile_section("run_visual_pipeline")
def run_visual_pipeline(
    file_path: Path,
    results: List[CheckResult],
    visual_pages: List[VisualPage],
    assets_root: Path,
) -> Optional[int]:
    suffix = file_path.suffix.lower()
    slug = slugify_for_path(file_path)
    pdf_dir = assets_root / "pdf"
    image_dir = assets_root / "png" / slug
    pdf_dir.mkdir(parents=True, exist_ok=True)


    if suffix in {".xls", ".xlsx", ".xlsm"}:
        sheet_pdfs, skipped_sheets, sheet_err = convert_excel_to_sheet_pdfs(file_path, pdf_dir, slug)
        if sheet_err:
            fallback_rendered, fallback_err = render_excel_sheets_to_png_fallback(file_path, image_dir)
            if fallback_err:
                # 画像生成失敗時の詳細原因をエクセル出力にも明記
                add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "ERROR", f"{sheet_err} / {fallback_err}", "Office環境または必要ライブラリを確認して再実行。")
                # フォールバック失敗時もvisual_pagesにエラー画像を追加（ダミー画像やエラー内容を記載した画像を生成してもよい）
                return None

            for page_no, (sheet_name, image_path) in enumerate(fallback_rendered, start=1):
                visual_pages.append(VisualPage(str(file_path), page_no, str(image_path), sheet_name=sheet_name))

            add_result(
                results,
                file_path,
                "CommonVisual",
                "V1-WARN",
                "共通PDF出力・PNG化（フォールバック）",
                "WARN",
                f"{sheet_err} のため、Excel内容の簡易PNG化で代替しました。({len(fallback_rendered)}枚)",
                "pywin32を導入するとシート別PDF経由の高精度画像化に戻ります。",
            )
            return len(fallback_rendered)

        if skipped_sheets:
            add_result(
                results,
                file_path,
                "CommonVisual",
                "V1-WARN",
                "共通PDF出力・PNG化（シート別）",
                "WARN",
                f"印刷対象なし等でシート別PDF化をスキップ: {', '.join(skipped_sheets[:6])}"
                + (f" ほか{len(skipped_sheets)-6}件" if len(skipped_sheets) > 6 else ""),
                "対象シートの印刷範囲・印刷設定を確認してください。",
            )

        total_pages = 0
        global_page_no = 1
        for sheet_name, sheet_pdf in sheet_pdfs:
            sheet_image_dir = image_dir / sanitize_filename_for_path(sheet_name)
            images, image_err = render_pdf_to_pngs(sheet_pdf, sheet_image_dir)
            if image_err:
                add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "ERROR", image_err, "必要ライブラリ導入後に再実行。")
                return None
            for image_path in images:
                visual_pages.append(VisualPage(str(file_path), global_page_no, str(image_path), sheet_name=sheet_name))
                global_page_no += 1
                total_pages += 1

        add_result(
            results,
            file_path,
            "CommonVisual",
            "V1",
            "共通PDF出力・PNG化",
            "PASS",
            f"{total_pages}ページをPNG化。 / 指摘対象ページ：P1?P{total_pages}",
            "別シートで罫線切れ/表不完全/ページ数を目検。",
        )
        return total_pages


    # PPT/PPTX画像化処理（PowerPoint COM経由で高精度PNG化）
    if suffix in {".ppt", ".pptx"}:
        # PPT/PPTX: ファイル名・拡張子をimage_previewとresultsで統一（.pptに）
        base_name = file_path.stem
        unified_file_path = str(file_path.with_suffix('.ppt'))
        pdf_path = pdf_dir / f"{slug}.pdf"
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "ERROR", "pywin32 が未インストール。", "pip install pywin32")
            return None
        images = []
        image_err = None
        pptx_exported = False
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            pres = powerpoint.Presentations.Open(str(file_path), WithWindow=False)
            try:
                pres.ExportAsFixedFormat(str(pdf_path), 2)  # 2=PDF
                pptx_exported = True
            except Exception as exc_export:
                slide_img_dir = image_dir / "slide_jpeg"
                slide_img_dir.mkdir(parents=True, exist_ok=True)
                for idx, slide in enumerate(pres.Slides, start=1):
                    jpeg_path = str(slide_img_dir / f"slide{idx}.jpg")
                    slide.Export(jpeg_path, "JPG")
                    images.append(jpeg_path)
                add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "WARN", f"ExportAsFixedFormat不可のためSlide.ExportでJPEG化({len(images)}枚)で代替", "PowerPoint/pywin32環境を確認。")
            pres.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
        except Exception as exc:
            add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "WARN", f"PowerPointエクスポート失敗: {exc}。python-pptxで簡易PNG化を試行します。", "PowerPointインストール・ファイル状態確認。")
        if pptx_exported and pdf_path.exists():
            images, image_err = render_pdf_to_pngs(pdf_path, image_dir)
            if image_err:
                add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "ERROR", image_err, "必要ライブラリ導入後に再実行。")
                return None
        elif not images:
            try:
                from pptx import Presentation
                from PIL import Image, ImageDraw, ImageFont
                prs = Presentation(str(file_path))
                image_dir.mkdir(parents=True, exist_ok=True)
                max_pages = 0
                for idx, slide in enumerate(prs.slides, start=1):
                    width = prs.slide_width // 9525
                    height = prs.slide_height // 9525
                    img = Image.new("RGB", (width, height), "white")
                    draw = ImageDraw.Draw(img)
                    if slide.background and hasattr(slide.background, 'fill') and slide.background.fill.type is not None:
                        fill = slide.background.fill
                        if fill.type == 1 and hasattr(fill, 'foreground_color'):
                            rgb = fill.foreground_color.rgb
                            if rgb:
                                img.paste((rgb[0], rgb[1], rgb[2]), [0, 0, width, height])
                    for shape in slide.shapes:
                        if hasattr(shape, 'fill') and shape.fill.type is not None:
                            fill = shape.fill
                            if fill.type == 1 and hasattr(fill, 'foreground_color'):
                                rgb = fill.foreground_color.rgb
                                if rgb:
                                    left = int(shape.left // 9525)
                                    top = int(shape.top // 9525)
                                    w = int(shape.width // 9525)
                                    h = int(shape.height // 9525)
                                    draw.rectangle([left, top, left + w, top + h], fill=(rgb[0], rgb[1], rgb[2]))
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            left = int(shape.left // 9525)
                            top = int(shape.top // 9525)
                            font_size = 20
                            try:
                                font = ImageFont.truetype("meiryo.ttc", font_size)
                            except Exception:
                                font = ImageFont.load_default()
                            text = shape.text
                            draw.text((left, top), text, fill="black", font=font)
                    out_path = image_dir / f"slide{idx}.png"
                    img.save(out_path)
                    images.append(str(out_path))
                    max_pages += 1
                add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "WARN", f"PowerPoint PDFエクスポート不可のためpython-pptxで簡易PNG化({max_pages}枚)で代替(背景・図形対応)", "PowerPoint/pywin32環境を確認。")
            except Exception as exc2:
                add_result(results, unified_file_path, "PPT", "V1", "PPT画像化", "ERROR", f"python-pptx簡易PNG化も失敗: {exc2}", "pptxファイル/ライブラリ確認。")
                return None
        total_pages = 0
        for idx, image_path in enumerate(images, start=1):
            visual_pages.append(VisualPage(unified_file_path, idx, str(image_path), sheet_name=f"slide{idx}"))
            total_pages += 1
        # 画像化結果
        add_result(
            results,
            unified_file_path,
            "PPT",
            "V1",
            "PPT画像化",
            "PASS",
            f"{total_pages}スライドをPNG化（PowerPointエクスポート）。 / 指摘対象ページ：P1?P{total_pages}",
            "スライド内容を目検。",
        )
        # VSDと同様の内容チェック（共通PDF出力・PNG化）もresultsに追加
        add_result(
            results,
            unified_file_path,
            "CommonVisual",
            "V1",
            "共通PDF出力・PNG化",
            "PASS",
            f"{total_pages}ページをPNG化。 / 指摘対象ページ：P1?P{total_pages}",
            "別シートで罫線切れ/表不完全/ページ数を目検。",
        )
        return total_pages

    # --- VISIOファイル対応（Aspose/Visio COM/LibreOffice の順で利用） ---
    if suffix in {".vsd", ".vsdx"}:
        # VSD/VSDX: ファイル名・拡張子をimage_previewとresultsで統一
        base_name = file_path.stem
        unified_file_path = str(file_path.with_suffix('.vsd')) if suffix == '.vsd' else str(file_path.with_suffix('.vsdx'))
        pdf_path = pdf_dir / f"{slug}.pdf"
        images, image_err, backend_name = render_visio_with_pdf_backend(file_path, pdf_path, image_dir)
        if images:
            total_pages = 0
            for idx, image_path in enumerate(images, start=1):
                visual_pages.append(VisualPage(unified_file_path, idx, str(image_path), sheet_name=f"page{idx}"))
                total_pages += 1
            add_result(
                results,
                unified_file_path,
                "CommonVisual",
                "VISO",
                "VISIO→PNG変換",
                "PASS",
                f"{total_pages}ページをPNG化({backend_name}→PDF)。 / 指摘対象ページ：P1?P{total_pages}",
                "ページ内容を目検。",
            )
            return total_pages

        import subprocess
        soffice_path = find_soffice_path()
        if soffice_path is not None:
            try:
                image_dir.mkdir(parents=True, exist_ok=True)
                cmd = [soffice_path, "--headless", "--convert-to", "png", str(file_path), "--outdir", str(image_dir)]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                base = file_path.stem
                png_files = list(image_dir.glob(f"{base}*.png"))
                if not png_files:
                    detail = f"LibreOfficeでPNG変換失敗: {result.stderr}".strip()
                    if image_err:
                        detail = f"{image_err} / {detail}"
                    add_result(results, file_path, "CommonVisual", "VISO", "VISIO→PNG変換", "ERROR", detail, "ファイルや環境を確認してください。")
                    return None
                for idx, image_path in enumerate(sorted(png_files), 1):
                    visual_pages.append(VisualPage(str(file_path), idx, str(image_path), sheet_name=None))
                detail = f"{len(png_files)}ページをPNG化(LibreOffice)。"
                if image_err:
                    detail = f"Aspose/COMは利用不可のため {detail}"
                add_result(results, file_path, "CommonVisual", "VISO", "VISIO→PNG変換", "PASS", detail, "画像内容を確認してください。")
                return len(png_files)
            except Exception as exc:
                detail = f"LibreOffice変換失敗: {exc}"
                if image_err:
                    detail = f"{image_err} / {detail}"
                add_result(results, file_path, "CommonVisual", "VISO", "VISIO→PNG変換", "ERROR", detail, "ファイルや環境を確認してください。")
                return None

        detail = "LibreOfficeが見つかりません。"
        if image_err:
            detail = f"{image_err} / {detail}"
        add_result(
            results,
            file_path,
            "CommonVisual",
            "VISO",
            "VISIO→PNG変換",
            "MANUAL",
            f"{detail} Visio for Web(Graph API)による変換にはMicrosoftアカウント認証が必要です。",
            "Aspose.Diagramプラグイン、Visio、またはLibreOffice環境を確認してください。",
        )
        return None

    pdf_path: Optional[Path] = None
    if suffix == ".pdf":
        pdf_path = file_path
    elif suffix in {".doc", ".docx"}:
        pdf_path = pdf_dir / f"{slug}.pdf"
        err = convert_office_to_pdf(file_path, pdf_path)
        if err:
            add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "ERROR", err, "Office環境を確認して再実行。")
            return None
    else:
        add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "N/A", "この拡張子は共通PDF出力対象外。", "対応不要。")
        return None

    if pdf_path is None or (not pdf_path.exists()):
        add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "ERROR", "PDF生成に失敗。", "変換環境を確認して再実行。")
        return None

    images, image_err = render_pdf_to_pngs(pdf_path, image_dir)
    if image_err:
        add_result(results, file_path, "CommonVisual", "V1", "共通PDF出力・PNG化", "ERROR", image_err, "必要ライブラリ導入後に再実行。")
        return None

    for idx, image_path in enumerate(images, start=1):
        visual_pages.append(VisualPage(str(file_path), idx, str(image_path), sheet_name=""))

    add_result(
        results,
        file_path,
        "CommonVisual",
        "V1",
        "共通PDF出力・PNG化",
        "PASS",
        f"{len(images)}ページをPNG化。 / 指摘対象ページ：P1?P{len(images)}",
        "別シートで罫線切れ/表不完全/ページ数を目検。",
    )
    return len(images)


def append_max_page_detail(results: List[CheckResult], file_path: Path, page_count: Optional[int]) -> None:
    # Excelシート別PDF化スキップ時の「最大ページ数=XX」は、PDF化できたシートの最大ページ数を示します。COM無効や印刷対象なしの場合は0になります。
    if page_count is None:
        return
    target = str(file_path)
    for r in results:
        if r.file_path != target:
            continue
        if "最大ページ数=" in r.detail:
            continue
        r.detail = f"{r.detail} / 最大ページ数={page_count}"


def derive_applicability(status: str) -> str:
    if status == "N/A":
        return "NOT_APPLICABLE"
    return "APPLICABLE"


def derive_automation(status: str) -> str:
    if status == "N/A":
        return "NOT_SUPPORTED"
    if status == "MANUAL":
        return "MANUAL_ONLY"
    return "AUTO_CHECKED"


def display_status(status: str) -> str:
    if status == "PASS":
        return "〇"
    if status in {"FAIL", "ERROR"}:
        return "×"

    return "要チェック"

def display_file_path_for_log(file_path: str) -> str:
    """
    ログや進捗表示用にPPT/PPTX/VSD/VSDXの拡張子を統一して表示
    """
    p = Path(file_path)
    ext = p.suffix.lower()
    if ext in {'.ppt', '.pptx'}:
        return str(p.with_suffix('.ppt').name)
    elif ext in {'.vsd', '.vsdx'}:
        return str(p.with_suffix('.vsd').name)
    else:
        return p.name


def populate_image_preview_sheet(ws_images, visual_pages: Iterable[VisualPage], preview_pages_per_row: int = 6) -> None:

    ws_images.append(["file_name", "page_count", "image_labels..."])
    ws_images.column_dimensions["A"].width = 65
    ws_images.column_dimensions["B"].width = 10
    ws_images.column_dimensions["C"].width = 28

    # 成果物（file_path）ごとに画像を新しい行に横並びで配置
    import os
    from collections import defaultdict
    unique_visual_pages = {}
    for page in visual_pages:
        key = (page.file_path, page.page_no, page.sheet_name)
        if key in unique_visual_pages:
            continue
        unique_visual_pages[key] = page

    # file_pathごとにグループ化
    pages_by_file = defaultdict(list)
    for page in unique_visual_pages.values():
        pages_by_file[page.file_path].append(page)
    for file_path in pages_by_file:
        pages_by_file[file_path] = sorted(pages_by_file[file_path], key=lambda p: (p.page_no, p.sheet_name))

    MM_PER_PIXEL = 0.264583  # 1px = 0.264583mm (96dpi)
    V_MARGIN_MM = 1  # 上下1mm
    H_MARGIN_MM = 1  # 左右1mm
    V_MARGIN_PX = int(V_MARGIN_MM / MM_PER_PIXEL)
    H_MARGIN_PX = int(H_MARGIN_MM / MM_PER_PIXEL)
    start_col = 3  # 画像は3列目以降に配置
    base_row = 3
    cur_row = base_row
    FIXED_IMG_HEIGHT = 180  # 画像高さ(px)を固定（調整可）
    FILE_ROW_GAP = 3  # ファイル間の空白行数
    IMAGE_COL_GAP = 1  # 画像間の空白列数

    # resultsで使われているVSDファイル名一覧を作成
    vsd_names = set()
    for f in pages_by_file:
        p = Path(f)
        if p.suffix.lower() in {'.vsd', '.vsdx'}:
            vsd_names.add(p.with_suffix('.vsd').name)

    for file_idx, (file_path, pages) in enumerate(sorted(pages_by_file.items())):
        # PPT/PPTX/VSD/VSDXは拡張子を統一して表示
        p = Path(file_path)
        ext = p.suffix.lower()
        if ext in {'.ppt', '.pptx'}:
            display_name = p.with_suffix('.ppt').name
        elif ext in {'.vsd', '.vsdx'}:
            display_name = p.with_suffix('.vsd').name
        elif ext == '.pdf' and p.with_suffix('.vsd').name in vsd_names:
            display_name = p.with_suffix('.vsd').name
        else:
            display_name = p.name

        # 画像情報を準備
        img_objs = []
        img_paths = []
        valid_indices = []
        for idx, page in enumerate(pages):
            try:
                img_path = str(Path(page.image_path).resolve())
                if XLImage is not None and Path(img_path).exists() and os.path.getsize(img_path) > 0:
                    img = XLImage(img_path)
                    img_objs.append(img)
                    img_paths.append(img_path)
                    valid_indices.append(idx)
                else:
                    img_objs.append(None)
                    img_paths.append("")
            except Exception:
                img_objs.append(None)
                img_paths.append("")

        # 有効な画像がなければスキップ
        if not any(img_objs):
            cur_row += FILE_ROW_GAP
            continue

        # image_preview では絶対パスではなくファイル名を表示して視認性を上げる
        ws_images.cell(row=cur_row, column=1, value=display_name)
        # 総ページ数を2列目に表示
        ws_images.cell(row=cur_row, column=2, value=len(pages))

        # ファイル名の下に1行空ける
        cur_row += 1

        # 画像ラベル行
        col = start_col
        for idx, page in enumerate(pages):
            label = f"p{page.page_no}"
            if page.sheet_name:
                label += f" / {page.sheet_name}"
            ws_images.cell(row=cur_row, column=col, value=label)
            col += 1 + IMAGE_COL_GAP  # 画像間に1列空ける

        # 画像行
        cur_row += 1
        col = start_col
        for idx, page in enumerate(pages):
            img = img_objs[idx]
            img_path = img_paths[idx]
            if img is not None and img_path:
                try:
                    # 画像を固定高さにリサイズ
                    if getattr(img, "height", 0) and img.height > 0 and img.height != FIXED_IMG_HEIGHT:
                        ratio = FIXED_IMG_HEIGHT / img.height
                        img.width = int(img.width * ratio)
                        img.height = int(img.height * ratio)
                    ws_images.add_image(img, ws_images.cell(row=cur_row, column=col).coordinate)
                    col_letter = ws_images.cell(row=1, column=col).column_letter
                    ws_images.column_dimensions[col_letter].width = max(22, int(img.width / 7))
                except Exception as exc:
                    print(f"[ERROR] 画像貼付失敗: {img_path} / {exc}")
                    ws_images.cell(row=cur_row, column=col, value=f"画像貼付失敗: {exc}\n{img_path}")
            col += 1 + IMAGE_COL_GAP

        # 行の高さを固定
        ws_images.row_dimensions[cur_row].height = FIXED_IMG_HEIGHT + 2 * V_MARGIN_PX

        # ファイル間に空白行
        cur_row += FILE_ROW_GAP


def write_results_report_xlsx(results: Iterable[CheckResult], output_xlsx: Path) -> Optional[str]:
    if Workbook is None:
        return "openpyxl が未インストールのため結果xlsxを出力できません。"

    filtered_results = [r for r in results if not is_excluded_check(r.check_id, r.check_item)]
    output_xlsx.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws_results = wb.active
    ws_results.title = "results"
    ws_results.append([
        "file_path",
        "file_type",
        "check_id",
        "check_item",
        "applicability",
        "automation",
        "status",
        "detail",
        "suggested_action",
    ])

    for r in filtered_results:
        # PPT/PPTX/VSD/VSDXは拡張子を統一して出力
        p = Path(r.file_path)
        ext = p.suffix.lower()
        if ext in {'.ppt', '.pptx'}:
            file_path_disp = str(p.with_suffix('.ppt'))
        elif ext in {'.vsd', '.vsdx'}:
            file_path_disp = str(p.with_suffix('.vsd'))
        else:
            file_path_disp = r.file_path
        ws_results.append([
            file_path_disp,
            r.file_type,
            r.check_id,
            r.check_item,
            derive_applicability(r.status),
            derive_automation(r.status),
            display_status(r.status),
            r.detail,
            r.suggested_action,
        ])

    # 罫線・左上詰め
    from openpyxl.styles import Border, Side, Alignment
    thin = Side(border_style="thin", color="000000")
    border = Border(top=thin, left=thin, right=thin, bottom=thin)
    align = Alignment(horizontal="left", vertical="top", wrap_text=True)
    max_row = ws_results.max_row
    max_col = ws_results.max_column
    for row in ws_results.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col):
        for cell in row:
            cell.border = border
            cell.alignment = align

    # 左上詰め（A1から出力済みなのでOK）

    try:
        wb.save(output_xlsx)
        return None
    except Exception as exc:
        return f"results xlsx出力失敗: {exc}"


def write_image_preview_xlsx(
    visual_pages: Iterable[VisualPage],
    output_xlsx: Path,
    preview_pages_per_row: int = 6,
) -> Optional[str]:
    if Workbook is None:
        return "openpyxl が未インストールのためimage_preview xlsxを出力できません。"

    output_xlsx.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws_images = wb.active
    ws_images.title = "image_preview"
    populate_image_preview_sheet(ws_images, visual_pages, preview_pages_per_row)

    try:
        wb.save(output_xlsx)
        return None
    except Exception as exc:
        return f"image_preview xlsx出力失敗: {exc}"


def apply_sheet_table_style(ws, max_col: Optional[int] = None) -> None:
    from openpyxl.styles import Border, Side, Alignment

    thin = Side(border_style="thin", color="000000")
    border = Border(top=thin, left=thin, right=thin, bottom=thin)
    align = Alignment(horizontal="left", vertical="top", wrap_text=True)
    max_row = ws.max_row
    max_col = max_col or ws.max_column
    for row in ws.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col):
        for cell in row:
            cell.border = border
            cell.alignment = align


def populate_other_files_sheet(ws_other, other_files: Iterable[Path]) -> None:
    ws_other.title = "other_files"
    ws_other.append(["file_path", "file_name", "extension"])
    ws_other.column_dimensions["A"].width = 90
    ws_other.column_dimensions["B"].width = 40
    ws_other.column_dimensions["C"].width = 16

    for file_path in sorted({Path(p).resolve() for p in other_files}, key=lambda p: str(p).lower()):
        ws_other.append([str(file_path), file_path.name, file_path.suffix])

    apply_sheet_table_style(ws_other, max_col=3)


def populate_suggested_action_settings_sheet(
    ws_settings,
    settings: Iterable[SuggestedActionSetting],
) -> None:
    ws_settings.title = SUGGESTED_ACTION_SETTINGS_SHEET
    ws_settings.append([
        "enabled",
        "file_type",
        "check_id",
        "status",
        "excel_status",
        "check_item",
        "suggested_action",
    ])
    ws_settings.column_dimensions["A"].width = 10
    ws_settings.column_dimensions["B"].width = 16
    ws_settings.column_dimensions["C"].width = 12
    ws_settings.column_dimensions["D"].width = 12
    ws_settings.column_dimensions["E"].width = 12
    ws_settings.column_dimensions["F"].width = 48
    ws_settings.column_dimensions["G"].width = 48

    for setting in settings:
        excel_status = setting.excel_status or ("×" if setting.status in {"FAIL", "ERROR", "CROSS"} else "")
        ws_settings.append([
            "TRUE" if setting.enabled else "FALSE",
            setting.file_type,
            setting.check_id,
            setting.status,
            excel_status,
            setting.check_item,
            setting.suggested_action,
        ])

    apply_sheet_table_style(ws_settings, max_col=7)


def write_visual_report_xlsx(
    results: Iterable[CheckResult],
    visual_pages: Iterable[VisualPage],
    output_xlsx: Path,
    other_files: Optional[Iterable[Path]] = None,
    preview_pages_per_row: int = 6,
    suggested_action_settings: Optional[Iterable[SuggestedActionSetting]] = None,
) -> Optional[str]:
    if Workbook is None:
        return "openpyxl が未インストールのため画像シートxlsxを出力できません。"

    filtered_results = [r for r in results if not is_excluded_check(r.check_id, r.check_item)]
    output_xlsx.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws_results = wb.active
    ws_results.title = "results"
    ws_results.append([
        "file_path",
        "file_type",
        "check_id",
        "check_item",
        "applicability",
        "automation",
        "status",
        "detail",
        "suggested_action",
    ])

    for r in filtered_results:
        # PPT/PPTX/VSD/VSDXは拡張子を統一して出力
        p = Path(r.file_path)
        ext = p.suffix.lower()
        if ext in {'.ppt', '.pptx'}:
            file_path_disp = str(p.with_suffix('.ppt'))
        elif ext in {'.vsd', '.vsdx'}:
            file_path_disp = str(p.with_suffix('.vsd'))
        else:
            file_path_disp = r.file_path
        ws_results.append([
            file_path_disp,
            r.file_type,
            r.check_id,
            r.check_item,
            derive_applicability(r.status),
            derive_automation(r.status),
            display_status(r.status),
            r.detail,
            r.suggested_action,
        ])
    apply_sheet_table_style(ws_results)

    ws_images = wb.create_sheet("image_preview")
    populate_image_preview_sheet(ws_images, visual_pages, preview_pages_per_row=preview_pages_per_row)
    ws_other = wb.create_sheet("other_files")
    populate_other_files_sheet(ws_other, other_files or [])
    ws_settings = wb.create_sheet(SUGGESTED_ACTION_SETTINGS_SHEET)
    populate_suggested_action_settings_sheet(ws_settings, suggested_action_settings or [])

    try:
        wb.save(output_xlsx)
        return None
    except Exception as exc:
        return f"xlsx出力失敗: {exc}"


def coord_in_ranges(coord: str, ranges: List[Tuple[int, int, int, int]]) -> bool:
    if coordinate_to_tuple is None:
        return False
    row, col = coordinate_to_tuple(coord)
    for min_col, min_row, max_col, max_row in ranges:
        if min_row <= row <= max_row and min_col <= col <= max_col:
            return True
    return False


def parse_print_area_ranges(ws) -> List[Tuple[int, int, int, int]]:
    areas: List[str] = []
    if ws.print_area:
        if isinstance(ws.print_area, str):
            areas = [ws.print_area]
        else:
            areas = list(ws.print_area)
    parsed: List[Tuple[int, int, int, int]] = []
    for area in areas:
        cleaned = area.replace("$", "")
        if "!" in cleaned:
            cleaned = cleaned.split("!", 1)[1]
        for part in cleaned.split(","):
            part = part.strip()
            if ":" not in part:
                continue
            try:
                parsed.append(range_boundaries(part))
            except Exception:
                continue
    return parsed


def extract_excel_break_ids(break_container) -> List[int]:
    if break_container is None:
        return []
    ids: List[int] = []
    for br in getattr(break_container, "brk", []) or []:
        try:
            ids.append(int(getattr(br, "id")))
        except Exception:
            continue
    return sorted(set(ids))


def infer_excel_print_page(ws, row_idx: int, col_idx: int) -> str:
    row_break_ids = extract_excel_break_ids(getattr(ws, "row_breaks", None))
    col_break_ids = extract_excel_break_ids(getattr(ws, "col_breaks", None))
    row_page = 1 + sum(1 for bid in row_break_ids if row_idx > bid)
    col_page = 1 + sum(1 for bid in col_break_ids if col_idx > bid)
    if col_page == 1:
        return f"P{row_page}"
    return f"P{row_page}-{col_page}"


def infer_excel_print_page_from_breaks(row_break_ids: List[int], col_break_ids: List[int], row_idx: int, col_idx: int) -> str:
    row_page = 1 + bisect_right(row_break_ids, row_idx - 1)
    col_page = 1 + bisect_right(col_break_ids, col_idx - 1)
    if col_page == 1:
        return f"P{row_page}"
    return f"P{row_page}-{col_page}"


def iter_nonempty_cells(ws):
    cells = getattr(ws, "_cells", None)
    if isinstance(cells, dict):
        for (row_idx, col_idx), cell in sorted(cells.items()):
            if cell is None or cell.value is None:
                continue
            yield row_idx, col_idx, cell
        return

    max_row = ws.max_row or 1
    max_col = ws.max_column or 1
    for row_idx in range(1, max_row + 1):
        for col_idx in range(1, max_col + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            if cell.value is None:
                continue
            yield row_idx, col_idx, cell


def collect_formula_refs_excel(workbook) -> Set[str]:
    refs: Set[str] = set()
    for ws in workbook.worksheets:
        cells = getattr(ws, "_cells", None)
        if isinstance(cells, dict):
            iter_cells = cells.values()
        else:
            iter_cells = (cell for _, _, cell in iter_nonempty_cells(ws))
        for cell in iter_cells:
            value = getattr(cell, "value", None)
            if isinstance(value, str) and value.startswith("="):
                for ref in re.findall(r"\b([A-Za-z]{1,3}\d{1,7})\b", value):
                    refs.add(ref.upper())
    return refs


def is_excel_red_font(cell) -> bool:
    font = getattr(cell, "font", None)
    color = getattr(font, "color", None)
    if color is None:
        return False
    rgb = str(getattr(color, "rgb", "") or "").upper()
    if rgb.endswith("FF0000"):
        return True
    try:
        indexed = getattr(color, "indexed", None)
        if indexed is not None and int(indexed) == 10:
            return True
    except Exception:
        pass
    return False


def is_excel_shaded(cell) -> bool:
    fill = getattr(cell, "fill", None)
    pattern = str(getattr(fill, "patternType", "") or "").lower()
    return pattern not in {"", "none"}


def is_excel_yellow_fill(cell) -> bool:
    fill = getattr(cell, "fill", None)
    fg = getattr(fill, "fgColor", None)
    if fg is None:
        return False
    rgb = str(getattr(fg, "rgb", "") or "").upper()
    if rgb.endswith("FFFF00"):
        return True
    try:
        indexed = getattr(fg, "indexed", None)
        if indexed is not None and int(indexed) in {6, 13}:
            return True
    except Exception:
        pass
    return False


def is_excel_underlined(cell) -> bool:
    underline = getattr(getattr(cell, "font", None), "underline", None)
    return bool(underline)


def summarize_pages(page_labels: List[str], limit: int = 8) -> str:
    if not page_labels:
        return "指摘対象ページ：なし"
    ordered = list(dict.fromkeys(page_labels))
    shown = ordered[:limit]
    rest = len(ordered) - len(shown)
    if rest > 0:
        return "指摘対象ページ：" + ", ".join(shown) + f" ほか{rest}件"
    return "指摘対象ページ：" + ", ".join(shown)

def summarize_locations(locations: List[str], limit: int = 8) -> str:

    if not locations:

        return ""

    shown = locations[:limit]

    rest = len(locations) - len(shown)

    if rest > 0:

        return ", ".join(shown) + f" ほか{rest}件"

    return ", ".join(shown)



def _to_ascii_digits(text: str) -> str:
    trans = str.maketrans("０１２３４５６７８９", "0123456789")
    return (text or "").translate(trans)


def _nth_weekday_of_month(year: int, month: int, weekday: int, nth: int) -> date:
    first = date(year, month, 1)
    shift = (weekday - first.weekday()) % 7
    return date(year, month, 1 + shift + (nth - 1) * 7)


def _vernal_equinox_day(year: int) -> int:
    return int(20.8431 + 0.242194 * (year - 1980) - ((year - 1980) // 4))


def _autumn_equinox_day(year: int) -> int:
    return int(23.2488 + 0.242194 * (year - 1980) - ((year - 1980) // 4))


def jp_holidays_for_year(year: int) -> Dict[date, str]:
    cached = _JP_HOLIDAY_CACHE.get(year)
    if cached is not None:
        return cached

    holidays: Dict[date, str] = {
        date(year, 1, 1): "元日",
        date(year, 2, 11): "建国記念の日",
        date(year, 2, 23): "天皇誕生日",
        date(year, 4, 29): "昭和の日",
        date(year, 5, 3): "憲法記念日",
        date(year, 5, 4): "みどりの日",
        date(year, 5, 5): "こどもの日",
        date(year, 8, 11): "山の日",
        date(year, 11, 3): "文化の日",
        date(year, 11, 23): "勤労感謝の日",
    }

    holidays[_nth_weekday_of_month(year, 1, 0, 2)] = "成人の日"
    holidays[_nth_weekday_of_month(year, 7, 0, 3)] = "海の日"
    holidays[_nth_weekday_of_month(year, 9, 0, 3)] = "敬老の日"
    holidays[_nth_weekday_of_month(year, 10, 0, 2)] = "スポーツの日"

    if 1948 <= year <= 2099:
        holidays[date(year, 3, _vernal_equinox_day(year))] = "春分の日"
        holidays[date(year, 9, _autumn_equinox_day(year))] = "秋分の日"

    start = date(year, 1, 1)
    end = date(year, 12, 31)
    cursor = start
    while cursor <= end:
        if cursor.weekday() < 5 and cursor not in holidays:
            if (cursor - timedelta(days=1)) in holidays and (cursor + timedelta(days=1)) in holidays:
                holidays[cursor] = "国民の休日"
        cursor += timedelta(days=1)

    for d in sorted(list(holidays.keys())):
        if d.weekday() != 6:
            continue
        sub = d
        while True:
            sub = sub + timedelta(days=1)
            if sub not in holidays:
                holidays[sub] = "振替休日"
                break

    _JP_HOLIDAY_CACHE[year] = holidays
    return holidays


def classify_calendar_day(target_day: date) -> Optional[str]:
    if target_day.weekday() >= 5:
        return "土日"
    holiday_name = jp_holidays_for_year(target_day.year).get(target_day)
    if holiday_name:
        return f"祝日({holiday_name})"
    return None


def _extract_date_candidates(text: str) -> List[Tuple[str, Optional[date], str]]:
    norm = _to_ascii_digits(text)
    candidates: List[Tuple[str, Optional[date], str]] = []

    patterns = [
        re.compile(r"(?<!\d)(20\d{2})[\./\-](\d{1,2})[\./\-](\d{1,2})(?!\d)"),
        re.compile(r"(?<!\d)(20\d{2})年\s*(\d{1,2})月\s*(\d{1,2})日"),
    ]

    seen = set()
    for pat in patterns:
        for m in pat.finditer(norm):
            raw = m.group(0)
            if raw in seen:
                continue
            seen.add(raw)
            y, mm, dd = int(m.group(1)), int(m.group(2)), int(m.group(3))
            try:
                parsed = date(y, mm, dd)
                candidates.append((raw, parsed, ""))
            except ValueError:
                candidates.append((raw, None, "不正日付"))

    return candidates


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "")).strip()


def resolve_cover_keywords(cover_keyword: Optional[str]) -> Tuple[List[str], bool]:
    if cover_keyword and str(cover_keyword).strip():
        raw = str(cover_keyword).replace("、", ",")
        kws = [normalize_text(x) for x in raw.split(",") if normalize_text(x)]
        return (kws or DEFAULT_COVER_KEYWORDS), True
    return DEFAULT_COVER_KEYWORDS, False


def evaluate_cover(first_page_text: str, cover_keyword: Optional[str]) -> Tuple[str, str, str]:
    text = first_page_text or ""
    keywords, user_specified = resolve_cover_keywords(cover_keyword)
    matched = [kw for kw in keywords if kw in text]

    if user_specified:
        if matched:
            return "PASS", f"指定キーワード一致: {', '.join(matched[:3])}。 / 指摘対象ページ：P1", "対応不要。"
        return "FAIL", f"指定キーワードを未検出: {', '.join(keywords[:4])}。 / 指摘対象ページ：P1", "規定表紙へ修正してください。"

    if matched:
        return "PASS", f"自動キーワード一致: {', '.join(matched[:3])}。 / 指摘対象ページ：P1", "対応不要。"
    return "WARN", "自動キーワードで規定表紙を特定できませんでした。 / 指摘対象ページ：P1", "必要に応じて表紙を確認してください。"


def safe_lower(text: str) -> str:
    return normalize_text(text).lower()


def is_probable_person_name(text: str) -> bool:
    s = normalize_text(text)
    if not s:
        return False
    if len(s) > 20:
        return False
    ng_words = ["PMO", "事務局", "チーム", "担当", "ベンダ", "関係者", "全員", "未定", "調整中", "会議", "課", "部", "室", "グループ"]
    if any(w.lower() in s.lower() for w in ng_words):
        return False
    if re.fullmatch(r"[A-Za-z]+(?:\s+[A-Za-z]+){0,2}", s):
        return True
    if re.fullmatch(r"[\u4e00-\u9fff]{2,6}", s):
        return True
    if re.fullmatch(r"[\u4e00-\u9fff]{1,4}\s+[\u4e00-\u9fff]{1,4}", s):
        return True
    if re.fullmatch(r"[\u30a1-\u30ff]{2,12}", s):
        return True
    return False


def parse_annex_numbers(text: str) -> Set[str]:
    norm = _to_ascii_digits(text)
    nums = set()
    for pat in [
        r"別紙\s*[-－]?(\d{1,3})",
        r"別添\s*[-－]?(\d{1,3})",
        r"別紙番号\s*[:：]?\s*(\d{1,3})",
    ]:
        for m in re.finditer(pat, norm, flags=re.IGNORECASE):
            nums.add(str(int(m.group(1))))
    return nums


def parse_wbs_numbers(text: str) -> Set[str]:
    norm = _to_ascii_digits(text)
    hits = set()
    for m in re.finditer(r"\b(\d+(?:[\.-]\d+){1,4})\b", norm):
        hits.add(m.group(1).replace("-", "."))
    for m in re.finditer(r"WBS\s*[:：]?\s*(\d+(?:[\.-]\d+){1,4})", norm, flags=re.IGNORECASE):
        hits.add(m.group(1).replace("-", "."))
    return hits


def parse_simple_sequence_numbers(text: str, label: str) -> List[int]:
    norm = _to_ascii_digits(text)
    nums = []
    for m in re.finditer(rf"{re.escape(label)}\s*([0-9]+)", norm):
        try:
            nums.append(int(m.group(1)))
        except Exception:
            pass
    return nums


def get_missing_sequence(nums: List[int]) -> List[int]:
    if not nums:
        return []
    uniq = sorted(set(n for n in nums if n > 0))
    if len(uniq) <= 1:
        return []
    return [n for n in range(uniq[0], uniq[-1] + 1) if n not in uniq]


def extract_excel_text_pages(wb) -> List[Tuple[str, str]]:
    text_pages: List[Tuple[str, str]] = []
    for ws in wb.worksheets:
        row_break_ids = extract_excel_break_ids(getattr(ws, "row_breaks", None))
        col_break_ids = extract_excel_break_ids(getattr(ws, "col_breaks", None))
        for row_idx, col_idx, cell in iter_nonempty_cells(ws):
            if isinstance(cell.value, str) and cell.value.strip():
                page_label = infer_excel_print_page_from_breaks(row_break_ids, col_break_ids, row_idx, col_idx)
                text_pages.append((cell.value, page_label))
    return text_pages


def extract_word_text_pages(doc, file_path: Path) -> List[Tuple[str, str]]:
    blocks: List[str] = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            blocks.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text and p.text.strip():
                        blocks.append(p.text)
    page_map = find_word_text_page_numbers(file_path, blocks[:120], max_snippets=120)
    out: List[Tuple[str, str]] = []
    for t in blocks:
        pages = page_map.get(normalize_word_snippet(t), [])
        out.append((t, f"P{pages[0]}" if pages else "特定不可(Wordレイアウト依存)"))
    return out


def extract_pdf_text_pages(reader) -> List[Tuple[str, str]]:
    text_pages: List[Tuple[str, str]] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            text_pages.append((page_text, f"P{page_index}"))
    return text_pages


def load_visio_diagram_via_aspose(file_path: Path):
    ok, msg = ensure_aspose_diagram_available(auto_install=False)
    if not ok:
        return None, f"Aspose.Diagram を利用できません(Python {sys.version_info.major}.{sys.version_info.minor}): {msg}"
    try:
        from aspose.diagram import Diagram
    except Exception as exc:
        return None, f"Aspose.Diagram を利用できません(Python {sys.version_info.major}.{sys.version_info.minor}): {exc}"

    try:
        return Diagram(str(file_path)), None
    except Exception as exc:
        return None, f"VISIO読込失敗: {exc}"


def iter_visio_shapes(shape_collection) -> Iterable[object]:
    count = getattr(shape_collection, "count", 0) or 0
    for idx in range(count):
        shape = shape_collection[idx]
        yield shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from iter_visio_shapes(child_shapes)


def extract_visio_text_pages(diagram) -> List[Tuple[str, str]]:
    text_pages: List[Tuple[str, str]] = []
    page_count = getattr(getattr(diagram, "pages", None), "count", 0) or 0
    for page_index in range(page_count):
        page = diagram.pages[page_index]
        page_label = f"P{page_index + 1}"
        page_name = str(getattr(page, "name", "") or "").strip()
        if page_name:
            text_pages.append((page_name, page_label))
        for shape in iter_visio_shapes(page.shapes):
            try:
                shape_text = str(shape.get_pure_text() or "").strip()
            except Exception:
                shape_text = ""
            if shape_text:
                text_pages.append((shape_text, page_label))
    return text_pages


def summarize_visio_custom_props(custom_props) -> List[str]:
    summaries: List[str] = []
    collection_items: List[object] = []
    count_attr = getattr(custom_props, "count", None)
    if isinstance(count_attr, int):
        for idx in range(count_attr):
            try:
                collection_items.append(custom_props[idx])
            except Exception:
                continue
    else:
        try:
            collection_items = list(custom_props)
        except Exception:
            collection_items = []

    for idx, item in enumerate(collection_items):
        name = normalize_text(str(getattr(item, "name", "") or "")) or f"custom_prop_{idx + 1}"
        prop_type = getattr(item, "prop_type", None)
        value_obj = getattr(item, "custom_value", None)
        value_text = ""
        value_kind = "empty"
        if value_obj is not None:
            string_value = getattr(value_obj, "value_string", None)
            if string_value is not None and str(string_value).strip():
                value_text = str(string_value).strip()
                value_kind = "string"
            else:
                date_value = getattr(value_obj, "value_date", None)
                if getattr(date_value, "year", 1) > 1:
                    value_text = str(date_value)
                    value_kind = "date"
                else:
                    number_value = getattr(value_obj, "value_number", None)
                    if number_value not in {None, 0, 0.0}:
                        value_text = str(number_value)
                        value_kind = "number"
                    else:
                        bool_value = getattr(value_obj, "value_bool", None)
                        if bool_value is True:
                            value_text = str(bool_value)
                            value_kind = "bool"
        if value_text:
            summaries.append(f"{name}[kind={value_kind}, prop_type={prop_type}]={value_text}")
        else:
            summaries.append(f"{name}[kind={value_kind}, prop_type={prop_type}]")
    return summaries


def check_ppt(file_path: Path, results: List[CheckResult], cover_keyword: Optional[str]) -> None:
    """PPT/PPTX basic checks. Detailed layout is confirmed through image_preview."""
    try:
        from pptx import Presentation
    except ImportError:
        add_result(results, file_path, "PPT", "P-ENV", "PPT解析ライブラリ", "ERROR", "python-pptx が未インストール。", "pip install python-pptx")
        return

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        add_result(results, file_path, "PPT", "P-OPEN", "PPTファイル読込", "ERROR", f"読込失敗: {exc}", "ファイル状態確認。")
        return

    props = prs.core_properties
    has_props = any([
        props.title, props.subject, props.keywords, props.category,
        props.author, props.last_modified_by, props.comments
    ])
    props_detail = (
        f"タイトル: {props.title}, 件名: {props.subject}, キーワード: {props.keywords}, "
        f"分類: {props.category}, 作成者: {props.author}, 管理者: {props.last_modified_by}, コメント: {props.comments}"
    )
    add_result(results, file_path, "PPT", "C1", "プロパティ情報削除", "FAIL" if has_props else "PASS", props_detail, "不要プロパティを削除。" if has_props else "対応不要。")

    slide_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))[:3000]
    c2_status, c2_detail, c2_action = evaluate_cover(slide_text, cover_keyword)
    add_result(results, file_path, "PPT", "C2", "表紙が規定のもの", c2_status, c2_detail, c2_action)

    add_result(results, file_path, "PPT", "C3", COMMON_MARGIN_RULE_TEXT, "N/A", "PPTはページ余白を安定取得できないため対象外。画像シートでレイアウトを確認してください。", "対応不要。")
    slide_count = len(prs.slides)
    add_result(results, file_path, "PPT", "C4", "ページ番号", "MANUAL", f"スライド数: {slide_count}。 / 指摘対象ページ：目検で確認", "連番整合を目視確認。")
    add_result(results, file_path, "PPT", "C5", "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）", "MANUAL", "別シートPNGで目検。 / 指摘対象ページ：目検で確認", "罫線切れ・表不完全・ページ数を確認。")


def check_visio(file_path: Path, results: List[CheckResult], cover_keyword: Optional[str]) -> None:
    diagram, err = load_visio_diagram_via_aspose(file_path)
    if diagram is None:
        add_result(results, file_path, "VISIO", "V-ENV", "VISIO解析ライブラリ", "WARN", err or "Aspose.Diagram を利用できません。", "Aspose.Diagramプラグイン、Visio、またはLibreOffice環境を確認してください。")
        return

    props = getattr(diagram, "document_props", None)
    custom_props = getattr(props, "custom_props", None) if props is not None else None
    has_props, props_detail = _build_property_detail(
        [
            ("タイトル", getattr(props, "title", None)),
            ("件名", getattr(props, "subject", None)),
            ("タグ", getattr(props, "keywords", None)),
            ("分類", getattr(props, "category", None)),
            ("作成者", getattr(props, "creator", None)),
            ("会社", getattr(props, "company", None)),
            ("管理者", getattr(props, "manager", None)),
            ("説明", getattr(props, "desc", None)),
        ]
    )
    add_result(results, file_path, "VISIO", "C1", "プロパティ情報削除", "FAIL" if has_props else "PASS", props_detail, "不要プロパティを削除。" if has_props else "対応不要。")

    text_pages = extract_visio_text_pages(diagram)
    first_page_text = "\n".join(text for text, page in text_pages if page == "P1")[:3000]
    c2_status, c2_detail, c2_action = evaluate_cover(first_page_text, cover_keyword)
    add_result(results, file_path, "VISIO", "C2", "表紙が規定のもの", c2_status, c2_detail, c2_action)

    add_result(results, file_path, "VISIO", "C3", COMMON_MARGIN_RULE_TEXT, "N/A", "余白チェックは対象外。", "対応不要。")

    page_count = getattr(getattr(diagram, "pages", None), "count", 0) or 0
    add_result(results, file_path, "VISIO", "C4", "ページ番号", "MANUAL", f"ページ数: {page_count}。 / 指摘対象ページ：目検で確認", "連番整合を目視確認。")
    add_result(results, file_path, "VISIO", "C5", "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）", "MANUAL", "別シートPNGで目検。 / 指摘対象ページ：目検で確認", "罫線切れ・表不完全・ページ数を確認。")

    add_result(results, file_path, "VISIO", "V1", "印刷範囲外記載チェック", "MANUAL", "VISIOの印刷範囲外判定は未対応。 / 指摘対象ページ：目検で確認", "PNG化結果を確認してください。")
    add_result(results, file_path, "VISIO", "V2", "不要ページチェック", "WARN" if page_count > 1 else "MANUAL", (f"ページ数: {page_count}。不要ページ有無を確認してください。 / 指摘対象ページ：目検で確認" if page_count > 1 else "単一ページです。不要ページ有無は目検で確認してください。 / 指摘対象ページ：目検で確認"), "不要ページがあれば削除してください。")
    add_result(results, file_path, "VISIO", "V3", "印刷プレビュー（見切れ/罫線欠け）", "MANUAL", "VISIOの見切れ/罫線欠けはPNGで目検。 / 指摘対象ページ：目検で確認", "PNG化結果を確認してください。")
    add_result(results, file_path, "VISIO", "V4", "見え消し（取り消し線/訂正線/非表示）残存", "MANUAL", "Aspose.Diagram での見え消し自動判定は未対応。 / 指摘対象ページ：目検で確認", "必要に応じて元ファイルを確認してください。")

    comment_like_hits: List[str] = []
    ref_error_pages: List[str] = []
    hyperlink_pages: List[str] = []
    for page_index in range(page_count):
        page = diagram.pages[page_index]
        page_label = f"P{page_index + 1}"
        for shape in iter_visio_shapes(page.shapes):
            try:
                shape_text = str(shape.get_pure_text() or "")
            except Exception:
                shape_text = ""
            if re.search(r"(?i)(comment|コメント|review)", shape_text):
                comment_like_hits.append(page_label)
            if re.search(r"(?i)(error!\s*reference source not found\.?|#ref!|#name\?|#value!|#div/0!|#num!|#null!|#n/a)", shape_text):
                ref_error_pages.append(page_label)
            try:
                hyperlink_count = getattr(getattr(shape, "hyperlinks", None), "count", 0) or 0
            except Exception:
                hyperlink_count = 0
            if hyperlink_count:
                hyperlink_pages.append(page_label)

    add_result(
        results,
        file_path,
        "VISIO",
        "V5",
        "コメント（吹き出し）残存",
        "WARN" if comment_like_hits else "MANUAL",
        (f"コメントらしきテキスト候補を検出。 / {summarize_pages(comment_like_hits)}" if comment_like_hits else "Aspose.Diagram ではコメント構造の安定抽出が未対応です。 / 指摘対象ページ：目検で確認"),
        "必要に応じて元ファイルを確認してください。",
    )
    add_result(results, file_path, "VISIO", "V6", "参照エラー残存", "FAIL" if ref_error_pages else "PASS", (f"参照エラー文字列を検出。 / {summarize_pages(ref_error_pages)}" if ref_error_pages else "参照エラー文字列なし。 / 指摘対象ページ：なし"), "参照元を修正してください。" if ref_error_pages else "対応不要。")





def run_language_consistency_checks(
    results: List[CheckResult],
    file_path: Path,
    file_type: str,
    text_pages: List[Tuple[str, str]],
) -> None:
    return





def mm_from_emu(value) -> Optional[float]:
    if value is None:
        return None
    try:
        return float(value) / EMU_PER_MM
    except Exception:
        return None


def normalize_word_snippet(text: str, max_len: int = 60) -> str:
    cleaned = re.sub(r"\s+", " ", (text or "")).strip()
    return cleaned[:max_len]


def find_word_text_page_numbers(
    doc_path: Path,
    snippets: List[str],
    max_snippets: int = 30,
    max_hits_per_snippet: int = 5,
    max_seconds: float = 8.0,
) -> Dict[str, List[int]]:
    if gencache is None:
        return {}

    uniq_snippets: List[str] = []
    seen = set()
    for snippet in snippets:
        s = normalize_word_snippet(snippet)
        if (not s) or (len(s) < 3) or (s in seen):
            continue
        seen.add(s)
        uniq_snippets.append(s)
    if not uniq_snippets:
        return {}

    try:
        cache_key = str(doc_path.resolve())
    except Exception:
        cache_key = str(doc_path)
    cache = _WORD_PAGE_MAP_CACHE.setdefault(cache_key, {})

    page_map: Dict[str, List[int]] = {}
    pending: List[str] = []
    for snippet in uniq_snippets[:max_snippets]:
        if snippet in cache:
            if cache[snippet]:
                page_map[snippet] = list(cache[snippet])
            continue
        pending.append(snippet)

    if not pending:
        return page_map

    word = None
    doc = None
    try:
        word = gencache.EnsureDispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(doc_path), ReadOnly=True)
        content_end = doc.Content.End
        search_start = time.perf_counter()

        for snippet in pending:
            if (time.perf_counter() - search_start) >= max_seconds:
                break
            rng = doc.Content
            pages = set()
            last_end = -1
            steps = 0
            while True:
                if (time.perf_counter() - search_start) >= max_seconds:
                    break
                if steps >= 200:
                    break
                try:
                    finder = rng.Find
                    finder.ClearFormatting()
                    finder.Text = snippet
                    finder.Forward = True
                    finder.Wrap = 0
                    finder.MatchWildcards = False
                    finder.MatchWholeWord = False
                    finder.MatchCase = False
                    executed = bool(finder.Execute())
                except Exception:
                    break
                if not executed:
                    break
                try:
                    page_no = int(rng.Information(3))
                except Exception:
                    page_no = None
                if page_no:
                    pages.add(page_no)
                if len(pages) >= max_hits_per_snippet:
                    break
                next_start = rng.End
                if next_start >= content_end:
                    break
                if next_start <= last_end:
                    break
                try:
                    rng.SetRange(next_start, content_end)
                except Exception:
                    break
                last_end = next_start
                steps += 1

            if pages:
                sorted_pages = sorted(pages)
                cache[snippet] = sorted_pages
                page_map[snippet] = sorted_pages
            else:
                cache[snippet] = []
    except Exception:
        return page_map
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass

    return page_map


def collect_pages_from_snippets(snippets: List[str], page_map: Dict[str, List[int]]) -> List[str]:
    pages = set()
    for snippet in snippets:
        s = normalize_word_snippet(snippet)
        for page_no in page_map.get(s, []):
            pages.add(f"P{page_no}")
    return sorted(pages)


def iter_word_runs(doc):
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            yield run
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        yield run


def color_is_pure_blue_word(color_obj) -> Optional[bool]:
    if color_obj is None:
        return None
    rgb = str(getattr(color_obj, "rgb", "") or "").upper()
    if not rgb:
        return None
    return rgb.endswith("0000FF")


def color_is_blue_word(color_obj) -> bool:
    if color_obj is None:
        return False
    rgb = str(getattr(color_obj, "rgb", "") or "").upper()
    if rgb.endswith("0000FF"):
        return True
    if rgb.endswith("000080") or rgb.endswith("3333FF") or rgb.endswith("0000CC"):
        return True
    return False


def color_is_pure_blue_excel(cell) -> Optional[bool]:
    color = getattr(getattr(cell, "font", None), "color", None)
    if color is None:
        return None
    rgb = str(getattr(color, "rgb", "") or "").upper()
    if not rgb:
        return None
    return rgb.endswith("0000FF")


def color_is_blue_excel(cell) -> bool:
    color = getattr(getattr(cell, "font", None), "color", None)
    if color is None:
        return False
    rgb = str(getattr(color, "rgb", "") or "").upper()
    if rgb.endswith("0000FF") or rgb.endswith("0000CC") or rgb.endswith("3333FF") or rgb.endswith("000080"):
        return True
    try:
        indexed = getattr(color, "indexed", None)
        if indexed is not None and int(indexed) in {5, 12}:
            return True
    except Exception:
        pass
    return False












def check_excel(file_path: Path, results: List[CheckResult]) -> None:

    if load_workbook is None:

        add_result(

            results,

            file_path,

            "Excel",

            "E-ENV",

            "Excel解析ライブラリ",

            "ERROR",

            "openpyxl が未インストールのため解析できません。",

            "pip install openpyxl を実行してください。",

        )

        return

 

    try:

        with warnings.catch_warnings():

            warnings.filterwarnings(

                "ignore",

                message=r"(?i)cannot parse header or footer so it will be ignored",

                category=UserWarning,

            )

            wb = load_workbook(file_path, data_only=False)

    except Exception as exc:

        add_result(

            results,

            file_path,

            "Excel",

            "E-OPEN",

            "Excelファイル読込",

            "ERROR",

            f"読込失敗: {exc}",

            "ファイル破損・パスワード保護の有無を確認してください。",

        )

        return

 

    all_formula_refs = collect_formula_refs_excel(wb)

 

    strike_cells = 0

    comment_cells = 0

    strike_locations: List[str] = []

    comment_locations: List[str] = []

    ref_error_locations: List[str] = []

    strike_pages: List[str] = []

    comment_pages: List[str] = []

    ref_error_pages: List[str] = []

    missing_print_title_rows_sheets: List[str] = []

 

    ref_error_tokens = {

        "#REF!",

        "#NAME?",

        "#VALUE!",

        "#DIV/0!",

        "#NUM!",

        "#NULL!",

        "#N/A",

    }

 

    for ws in wb.worksheets:

        print_title_rows = getattr(ws, "print_title_rows", None)

        if isinstance(print_title_rows, str) and print_title_rows.strip():

            pass

        else:

            missing_print_title_rows_sheets.append(ws.title)

 

        print_ranges = parse_print_area_ranges(ws)

        row_break_ids = extract_excel_break_ids(getattr(ws, "row_breaks", None))

        col_break_ids = extract_excel_break_ids(getattr(ws, "col_breaks", None))

        outside_cells: List[str] = []

        outside_locations: List[str] = []

        outside_referenced_locations: List[str] = []

        outside_pages: List[str] = []

        outside_referenced_pages: List[str] = []

 

        for row_idx, col_idx, cell in iter_nonempty_cells(ws):

            page_label = infer_excel_print_page_from_breaks(row_break_ids, col_break_ids, row_idx, col_idx)

            coord = f"{get_column_letter(col_idx)}{row_idx}"

            loc = f"{ws.title}!{coord}({page_label})"

 

            if bool(getattr(getattr(cell, "font", None), "strike", False)):

                strike_cells += 1

                strike_locations.append(loc)

                strike_pages.append(page_label)

            if getattr(cell, "comment", None) is not None:

                comment_cells += 1

                comment_locations.append(loc)

                comment_pages.append(page_label)

 

            cell_text = str(cell.value).upper() if isinstance(cell.value, str) else ""

            if (isinstance(cell.value, str) and cell.value.startswith("=") and "#REF!" in cell_text) or (

                cell_text in ref_error_tokens

            ):

                ref_error_locations.append(loc)

                ref_error_pages.append(page_label)

 

            if print_ranges and (not coord_in_ranges(coord, print_ranges)):

                outside_cells.append(coord)

                outside_locations.append(loc)

                outside_pages.append(page_label)

                if coord in all_formula_refs:

                    outside_referenced_locations.append(loc)

                    outside_referenced_pages.append(page_label)

 

        if not print_ranges:

            add_result(

                results,

                file_path,

                "Excel",

                "E1",

                f"印刷範囲外記載チェック [{ws.title}]",

                "WARN",

                "印刷範囲が未設定のため、自動判定できません。 / 指摘対象ページ：特定不可(印刷範囲未設定)",

                "必要なら印刷範囲を設定し、再チェックしてください。",

            )

            continue

 

        if outside_cells:

            add_result(

                results,

                file_path,

                "Excel",

                "E1",

                f"印刷範囲外記載チェック [{ws.title}]",

                "FAIL",

                (

                    f"印刷範囲外セル {len(outside_cells)} 件（うち参照あり {len(outside_referenced_locations)} 件）。"

                    f" / {summarize_pages(outside_pages)}"

                    f" / 該当セル: {summarize_locations(outside_locations)}"

                    + (

                        f" / 参照ありセル: {summarize_locations(outside_referenced_locations)} / {summarize_pages(outside_referenced_pages)}"

                        if outside_referenced_locations

                        else ""

                    )

                ),

                "必要なら印刷範囲拡大、不要かつ未参照なら削除、参照ありは残置を検討。",

            )

        else:

            add_result(

                results,

                file_path,

                "Excel",

                "E1",

                f"印刷範囲外記載チェック [{ws.title}]",

                "PASS",

                "印刷範囲外の記載は検出されませんでした。 / 指摘対象ページ：なし",

                "対応不要。",

            )

 

    all_sheet_names = [ws.title for ws in wb.worksheets]

    add_result(

        results,

        file_path,

        "Excel",

        "E2",

        "不要シートチェック",

        "MANUAL",

        f"対象シート: {', '.join(all_sheet_names)}",

        "不要シート有無は対象シート一覧をもとに手動確認。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "E3",

        "印刷プレビュー（見切れ/罫線欠け）",

        "MANUAL",

        "プレビュー見切れ・罫線欠けは自動判定が難しいため目視確認が必要です。",

        "Excelで印刷プレビューを開き、見切れや罫線欠けを補正。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "E4",

        "見え消し（取り消し線/訂正線）残存",

        "FAIL" if strike_cells > 0 else "PASS",

        (

            f"取り消し線セル: {strike_cells} 件。 / {summarize_pages(strike_pages)} / 該当セル: {summarize_locations(strike_locations)}"

            if strike_cells > 0

            else "取り消し線セル: 0 件。 / 指摘対象ページ：なし"

        ),

        "不要な取り消し線を解除。" if strike_cells > 0 else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "E5",

        "コメント（吹き出し）残存",

        "FAIL" if comment_cells > 0 else "PASS",

        (

            f"コメント付きセル: {comment_cells} 件。 / {summarize_pages(comment_pages)} / 該当セル: {summarize_locations(comment_locations)}"

            if comment_cells > 0

            else "コメント付きセル: 0 件。 / 指摘対象ページ：なし"

        ),

        "不要なコメント/メモを削除。" if comment_cells > 0 else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "E6",

        "参照エラー残存",

        "FAIL" if ref_error_locations else "PASS",

        (

            f"参照エラー/数式エラーセル: {len(ref_error_locations)} 件。 / {summarize_pages(ref_error_pages)} / 該当セル: {summarize_locations(ref_error_locations)}"

            if ref_error_locations

            else "参照エラー/数式エラーセルは検出されません。 / 指摘対象ページ：なし"

        ),

        "該当セルの数式参照先を修正。" if ref_error_locations else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "E7",

        "ページ設定のタイトル行",

        "FAIL" if missing_print_title_rows_sheets else "PASS",

        (

            "タイトル行が未設定のシートがあります: " + ", ".join(missing_print_title_rows_sheets)

            if missing_print_title_rows_sheets

            else "全シートでタイトル行が設定されています。"

        ),

        "ページ設定のタイトル行（先頭行）を設定してください。"

        if missing_print_title_rows_sheets

        else "対応不要。",

    )

 

    cp = wb.properties

    has_props = any(

        [

            cp.creator,

            cp.lastModifiedBy,

            cp.title,

            cp.subject,

            cp.keywords,

            cp.description,

            cp.category,

            cp.contentStatus,

        ]

    )

    add_result(

        results,

        file_path,

        "Excel",

        "C1",

        "プロパティ情報削除",

        "FAIL" if has_props else "PASS",

        "プロパティ情報が残っています。" if has_props else "主要プロパティは空です。",

        "ファイル情報のプロパティを削除。" if has_props else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "C2",

        "表紙が規定のもの",

        "MANUAL",

        "規定表紙判定は基準データが必要なため手動確認です。",

        "規定表紙テンプレートと照合。",

    )

 

    margin_msgs = []
    margin_fail = False
    for ws in wb.worksheets:
        pm = ws.page_margins
        top = inches_to_mm(getattr(pm, "top", None))
        bottom = inches_to_mm(getattr(pm, "bottom", None))
        left = inches_to_mm(getattr(pm, "left", None))
        right = inches_to_mm(getattr(pm, "right", None))
        is_landscape = _excel_is_landscape(ws)
        ok, msg = _judge_c3_margin(top, bottom, left, right, is_landscape)
        margin_msgs.append(f"{ws.title}({msg})")
        if not ok:
            margin_fail = True

    add_result(
        results,
        file_path,
        "Excel",
        "C3",
        COMMON_MARGIN_RULE_TEXT,
        "FAIL" if margin_fail else "PASS",
        "; ".join(margin_msgs),
        "ページ設定で余白を補正。" if margin_fail else "設定値確認済み。",
    )

    add_result(
        results,
        file_path,
        "Excel",
        "C4",

        "ページ番号",

        "MANUAL",

        "Excelヘッダ/フッタのページ番号形式は目視確認を推奨。最大ページ数は画像シート参照。",

        "印刷プレビューまたはページ設定で確認。",

    )

 

    add_result(

        results,

        file_path,

        "Excel",

        "C5",

        "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）",

        "MANUAL",

        "PDF出力後の品質確認は目視が必要です。",

        "PDFを開き、見切れ・罫線欠け・表サイズ・ページ番号を確認。",

    )

 

    try:

        wb.close()

    except Exception:

        pass




def check_word(file_path: Path, results: List[CheckResult], cover_keyword: Optional[str]) -> None:

    if Document is None:

        add_result(

            results,

            file_path,

            "Word",

            "W-ENV",

            "Word解析ライブラリ",

            "ERROR",

            "python-docx が未インストールのため解析できません。",

            "pip install python-docx を実行してください。",

        )

        return

 

    try:

        doc = Document(file_path)

    except Exception as exc:

        add_result(

            results,

            file_path,

            "Word",

            "W-OPEN",

            "Wordファイル読込",

            "ERROR",

            f"読込失敗: {exc}",

            "ファイル破損・保護設定を確認してください。",

        )

        return

 

    doc_xml = doc.element.xml

    settings_xml = doc.settings.element.xml if getattr(doc, "settings", None) is not None else ""

    revision_patterns = [r"<w:ins\\b", r"<w:del\\b", r"<w:moveFrom\\b", r"<w:moveTo\\b"]

    revision_count = sum(len(re.findall(p, doc_xml)) for p in revision_patterns)

    track_revisions_on = "w:trackRevisions" in settings_xml

    strike_count = 0

    double_strike_count = 0

    highlighted_runs = 0

    colored_runs = 0

    strike_snippets: List[str] = []

    highlight_snippets: List[str] = []

    colored_snippets: List[str] = []

    ref_error_snippets: List[str] = []

    for run in iter_word_runs(doc):

        run_text = normalize_word_snippet(getattr(run, "text", ""))

        if bool(getattr(run.font, "strike", False)):

            strike_count += 1

            if run_text:

                strike_snippets.append(run_text)

        if bool(getattr(run.font, "double_strike", False)):

            double_strike_count += 1

            if run_text:

                strike_snippets.append(run_text)

        if run.font.highlight_color is not None:

            highlighted_runs += 1

            if run_text:

                highlight_snippets.append(run_text)

        if run.font.color is not None and run.font.color.rgb is not None:

            colored_runs += 1

            if run_text:

                colored_snippets.append(run_text)

 

    has_track_changes = revision_count > 0 or track_revisions_on

    has_miekeshi = has_track_changes or (strike_count + double_strike_count) > 0

    add_result(

        results,

        file_path,

        "Word",

        "W1",

        "見え消し（変更履歴）残存",

        "FAIL" if has_miekeshi else "PASS",

        (

            f"変更履歴要素={revision_count}, trackRevisions={track_revisions_on}, 取り消し線Run={strike_count}, 二重取り消し線Run={double_strike_count}。 / 指摘対象ページ：特定不可(Wordレイアウト依存)"

            if has_miekeshi

            else "見え消し（変更履歴/取り消し線/二重取り消し線）は検出されません。 / 指摘対象ページ：なし"

        ),

        "最新版確認後、変更履歴を承諾/破棄し、取り消し線を解消。" if has_miekeshi else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Word",

        "W2",

        "マーカ残存",

        "FAIL" if highlighted_runs > 0 else "PASS",

        (

            f"ハイライト付きRun: {highlighted_runs} 件。 / 指摘対象ページ：特定不可(Wordレイアウト依存)"

            if highlighted_runs > 0

            else "ハイライト付きRun: 0 件。 / 指摘対象ページ：なし"

        ),

        "最新版確認後、不要マーカを削除。" if highlighted_runs > 0 else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "Word",

        "W3",

        "不要な文字色",

        "WARN" if colored_runs > 0 else "PASS",

        (

            f"明示的な文字色付きRun: {colored_runs} 件。 / 指摘対象ページ：特定不可(Wordレイアウト依存)"

            if colored_runs > 0

            else "明示的な文字色付きRun: 0 件。 / 指摘対象ページ：なし"

        ),

        "不要な色がないことを最終確認し、色付き箇所は標準色へ修正。",

    )

 

    comment_anchor_count = len(re.findall(r"<w:commentRangeStart\\b", doc_xml)) + len(

        re.findall(r"<w:commentReference\\b", doc_xml)

    )

    has_comment_part = any("/comments.xml" in str(p.partname) for p in doc.part.package.parts)

    has_comments = comment_anchor_count > 0 or has_comment_part

    add_result(

        results,

        file_path,

        "Word",

        "W4",

        "コメント（吹き出し）残存",

        "FAIL" if has_comments else "PASS",

        (

            f"コメント参照を検出（anchor={comment_anchor_count}, comments.xml={has_comment_part}）。 / 指摘対象ページ：特定不可(Wordレイアウト依存)"

            if has_comments

            else "コメント（吹き出し）は検出されません。 / 指摘対象ページ：なし"

        ),

        "不要なコメントを削除し、最終版で吹き出し非表示を確認。" if has_comments else "対応不要。",

    )

 

    word_text_blocks = [p.text for p in doc.paragraphs if p.text]

    for table in doc.tables:

        for row in table.rows:

            for cell in row.cells:

                for p in cell.paragraphs:

                    if p.text:

                        word_text_blocks.append(p.text)

    ref_error_pattern = re.compile(r"(?i)(error!\s*reference source not found\.?|#ref!)")

    word_ref_hits = 0

    for t in word_text_blocks:

        if ref_error_pattern.search(t):

            word_ref_hits += 1

            ref_error_snippets.append(normalize_word_snippet(t))

 

    all_lookup_snippets = strike_snippets + highlight_snippets + colored_snippets + ref_error_snippets

    word_page_map = find_word_text_page_numbers(file_path, all_lookup_snippets)

    w1_pages = collect_pages_from_snippets(strike_snippets, word_page_map)

    w2_pages = collect_pages_from_snippets(highlight_snippets, word_page_map)

    w3_pages = collect_pages_from_snippets(colored_snippets, word_page_map)

    w5_pages = collect_pages_from_snippets(ref_error_snippets, word_page_map)

    add_result(

        results,

        file_path,

        "Word",

        "W5",

        "参照エラー残存",

        "FAIL" if word_ref_hits > 0 else "PASS",

        (

            f"参照エラー文字列を {word_ref_hits} 件検出。 / 指摘対象ページ：特定不可(Wordレイアウト依存)"

            if word_ref_hits > 0

            else "参照エラー文字列は検出されません。 / 指摘対象ページ：なし"

        ),

        "参照元(相互参照・図表番号・ブックマーク)を再設定。" if word_ref_hits > 0 else "対応不要。",

    )

 

    cp = doc.core_properties

    has_props = any(

        [

            cp.author,

            cp.last_modified_by,

            cp.title,

            cp.subject,

            cp.keywords,

            cp.comments,

            cp.category,

        ]

    )

    add_result(

        results,

        file_path,

        "Word",

        "C1",

        "プロパティ情報削除",

        "FAIL" if has_props else "PASS",

        "プロパティ情報が残っています。" if has_props else "主要プロパティは空です。",

        "ファイル情報のプロパティを削除。" if has_props else "対応不要。",

    )

 

    if cover_keyword:

        first_page_text = "\n".join(p.text for p in doc.paragraphs[:20])

        cover_ok = cover_keyword in first_page_text

        add_result(

            results,

            file_path,

            "Word",

            "C2",

            "表紙が規定のもの",

            "PASS" if cover_ok else "FAIL",

            f"表紙キーワード '{cover_keyword}' を " + ("検出。" if cover_ok else "未検出。"),

            "規定の表紙へ差し替え。" if not cover_ok else "対応不要。",

        )

    else:

        add_result(

            results,

            file_path,

            "Word",

            "C2",

            "表紙が規定のもの",

            "MANUAL",

            "cover keyword 未指定のため手動確認です。",

            "--cover-keyword を指定すると半自動確認できます。",

        )

 

    margin_fail_msgs = []
    margin_ok_msgs = []
    for idx, section in enumerate(doc.sections, start=1):
        top = mm_from_emu(section.top_margin)
        bottom = mm_from_emu(section.bottom_margin)
        left = mm_from_emu(section.left_margin)
        right = mm_from_emu(section.right_margin)
        is_landscape = _word_is_landscape(section)
        ok, msg = _judge_c3_margin(top, bottom, left, right, is_landscape)
        full_msg = f"sec{idx}({msg})"
        if ok:
            margin_ok_msgs.append(full_msg)
        else:
            margin_fail_msgs.append(full_msg)


    add_result(
        results,
        file_path,
        "Word",
        "C3",
        COMMON_MARGIN_RULE_TEXT,
        "FAIL" if margin_fail_msgs else "PASS",
        "; ".join(margin_fail_msgs if margin_fail_msgs else margin_ok_msgs),
        "ページ設定で余白を補正。" if margin_fail_msgs else "対応不要。",
    )

    # ページ番号フィールド検出ロジックを追加
    try:
        header_footer_xml = "\n".join(sec.header._element.xml + sec.footer._element.xml for sec in doc.sections)
        has_page_number_field = ("PAGE" in header_footer_xml) or ("w:pgNum" in header_footer_xml)
    except Exception:
        has_page_number_field = False

    add_result(
        results,
        file_path,
        "Word",
        "C4",
        "ページ番号",
        "PASS" if has_page_number_field else "WARN",
        "ページ番号フィールドを検出。" if has_page_number_field else "ページ番号フィールドを検出できません。最大ページ数は画像シート参照。",
        "ヘッダ/フッタにページ番号を設定。" if not has_page_number_field else "対応不要。",
    )

 

    add_result(

        results,

        file_path,

        "Word",

        "C5",

        "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）",

        "MANUAL",

        "PDF出力後の品質確認は目視が必要です。",

        "PDFを開き、見切れ・罫線欠け・表サイズ・ページ番号を確認。",

    )

 

    package = getattr(getattr(doc, "part", None), "package", None)

    close_fn = getattr(package, "close", None)

    if callable(close_fn):

        try:

            close_fn()

        except Exception:

            pass




def check_pdf(file_path: Path, results: List[CheckResult], cover_keyword: Optional[str]) -> None:

    if PdfReader is None:

        add_result(

            results,

            file_path,

            "PDF",

            "P-ENV",

            "PDF解析ライブラリ",

            "ERROR",

            "pypdf が未インストールのため解析できません。",

            "pip install pypdf を実行してください。"

        )

        return

 

    try:

        reader = PdfReader(str(file_path))

    except Exception as exc:

        add_result(

            results,

            file_path,

            "PDF",

            "P-OPEN",

            "PDFファイル読込",

            "ERROR",

            f"読込失敗: {exc}",

            "ファイル破損・パスワード・権限制限を確認してください。"

        )

        return

 

    metadata = reader.metadata or {}

    meaningful_meta = []

    for key in ["/Author", "/Creator", "/Producer", "/Title", "/Subject", "/Keywords"]:

        val = metadata.get(key)

        if val:

            meaningful_meta.append(f"{key}={val}")

 

    add_result(

        results,

        file_path,

        "PDF",

        "C1",

        "プロパティ情報削除",

        "FAIL" if meaningful_meta else "PASS",

        "; ".join(meaningful_meta) if meaningful_meta else "主要メタデータは空です。",

        "PDF作成時にメタデータ削除して再出力。" if meaningful_meta else "対応不要。"

    )

 

    if cover_keyword and reader.pages:

        first_text = (reader.pages[0].extract_text() or "")[:3000]

        cover_ok = cover_keyword in first_text

        add_result(

            results,

            file_path,

            "PDF",

            "C2",

            "表紙が規定のもの",

            "PASS" if cover_ok else "FAIL",

            f"表紙キーワード '{cover_keyword}' を " + ("検出。" if cover_ok else "未検出。"),

            "規定の表紙へ差し替え。" if not cover_ok else "対応不要。"

        )

    else:

        add_result(

            results,

            file_path,

            "PDF",

            "C2",

            "表紙が規定のもの",

            "MANUAL",

            "cover keyword 未指定または1ページ目取得不可のため手動確認です。",

            "--cover-keyword を指定すると半自動確認できます。",

        )

 

    add_result(

        results,

        file_path,

        "PDF",

        "C3",

        COMMON_MARGIN_RULE_TEXT,

        "MANUAL",

        "PDFはレイアウト固定で余白の厳密自動判定が難しいため手動確認です（基準: 上20/下20/左30/右20mm以上）。",

        "表示倍率100%で余白とパンチ穴干渉を確認。",

    )

 

    page_count = len(reader.pages)

    strikeout_annots = 0

    comment_annots = 0

    other_markup_annots = 0

    strike_pages: List[int] = []

    comment_pages: List[int] = []

    markup_pages: List[int] = []

    ref_error_pages: List[int] = []

    ref_error_pattern = re.compile(r"(?i)(error!\s*reference source not found\.?|#ref!)")

    for page_index, page in enumerate(reader.pages, start=1):

        try:

            page_text = page.extract_text() or ""

        except Exception:

            page_text = ""

        if ref_error_pattern.search(page_text):

            ref_error_pages.append(page_index)

 

        annots = page.get("/Annots")

        if not annots:

            continue

        for annot_ref in annots:

            try:

                annot = annot_ref.get_object()

            except Exception:

                continue

            subtype = annot.get("/Subtype")

            if subtype == "/StrikeOut":

                strikeout_annots += 1

                strike_pages.append(page_index)

            elif subtype in {"/Text", "/FreeText", "/Popup"}:

                comment_annots += 1

                comment_pages.append(page_index)

            elif subtype in {"/Underline", "/Squiggly", "/Highlight", "/Caret"}:

                other_markup_annots += 1

                markup_pages.append(page_index)

 

    add_result(

        results,

        file_path,

        "PDF",

        "P1",

        "見え消し（注釈ベース）残存",

        "FAIL" if (strikeout_annots + other_markup_annots) > 0 else "PASS",

        (

            f"注釈検出: StrikeOut={strikeout_annots}, その他マークアップ={other_markup_annots}。 / 指摘対象ページ：{', '.join(f'P{p}' for p in sorted(set(strike_pages + markup_pages)))}"

            if (strikeout_annots + other_markup_annots) > 0

            else "注釈ベースの見え消しは検出されません。 / 指摘対象ページ：なし"

        ),

        "不要な注釈を削除。テキストに焼き込まれた取り消し線は目視確認。"

        if (strikeout_annots + other_markup_annots) > 0

        else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "PDF",

        "P2",

        "コメント（吹き出し）残存",

        "FAIL" if comment_annots > 0 else "PASS",

        (

            f"コメント系注釈: {comment_annots} 件。 / 指摘対象ページ：{', '.join(f'P{p}' for p in sorted(set(comment_pages)))}"

            if comment_annots > 0

            else "コメント系注釈: 0 件。 / 指摘対象ページ：なし"

        ),

        "不要なコメント注釈を削除。" if comment_annots > 0 else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "PDF",

        "P3",

        "参照エラー残存",

        "FAIL" if ref_error_pages else "PASS",

        (

            f"参照エラー文字列を検出。 / 指摘対象ページ：{', '.join(f'P{p}' for p in sorted(set(ref_error_pages)))}"

            if ref_error_pages

            else "参照エラー文字列は検出されません。 / 指摘対象ページ：なし"

        ),

        "参照元を修正し、PDF再出力。" if ref_error_pages else "対応不要。",

    )

 

    add_result(

        results,

        file_path,

        "PDF",

        "C4",

        "ページ番号",

        "MANUAL",

        f"ページ数: {page_count}。抽出品質の都合でページ番号整合は手動確認を推奨。",

        "全ページのページ番号連番を確認。",

    )

 

    add_result(

        results,

        file_path,

        "PDF",

        "C5",

        "PDF出力結果確認（見切れ/罫線/表サイズ/ページ番号）",

        "MANUAL",

        "PDF品質確認の総合判定は目視が必要です。詳細は C5.1?C5.4 を参照。",

        "C5.1?C5.4 を確認して総合判断。",

    )

 

    add_result(

        results,

        file_path,

        "PDF",

        "C5.1",

        "PDF見切れ",

        "MANUAL",

        "見切れは表示環境依存のため目視確認が必要です。",

        "各ページ端部の文字欠けを確認。",

    )

    add_result(

        results,

        file_path,

        "PDF",

        "C5.2",

        "PDF罫線欠け",

        "MANUAL",

        "罫線欠けの判定は目視確認が必要です。",

        "細線・表罫線の欠けを確認。",

    )

    add_result(

        results,

        file_path,

        "PDF",

        "C5.3",

        "PDF表サイズ極小",

        "MANUAL",

        "表サイズの妥当性は文書意図依存のため目視確認です。",

        "表が過度に縮小されていないか確認。",

    )

    add_result(

        results,

        file_path,

        "PDF",

        "C5.4",

        "PDFページ番号正当性",

        "MANUAL",

        "ページ番号表示と総ページ整合は目視確認です。",

        "ページ番号の欠番・重複を確認。",

    )

 

    reader_stream = getattr(reader, "stream", None)

    if reader_stream is not None:

        try:

            reader_stream.close()

        except Exception:

            pass



def check_file(file_path: Path, results: List[CheckResult], cover_keyword: Optional[str]) -> None:
    suffix = file_path.suffix.lower()
    file_type = ""
    if suffix in {".xlsx", ".xlsm"}:
        file_type = "Excel"
        check_excel(file_path, results)
    elif suffix == ".docx":
        file_type = "Word"
        check_word(file_path, results, cover_keyword)
    elif suffix == ".pdf":
        file_type = "PDF"
        check_pdf(file_path, results, cover_keyword)
    elif suffix in {".ppt", ".pptx"}:
        file_type = "PPT"
        check_ppt(file_path, results, cover_keyword)
    elif suffix in {".xls", ".doc"}:
        file_type = "LegacyOffice"
        add_result(results, file_path, "LegacyOffice", "L1", "旧形式ファイル", "WARN", "旧形式(.xls/.doc)は詳細解析対象外。 / 指摘対象ページ：変換後に確認", "可能なら .xlsx/.docx へ変換。")

    # V5基準: 実際にチェックした結果のみ出力する。未対応/N/Aの補完は行わない。


def find_target_files(root: Path, exclude_paths: Optional[Set[Path]] = None) -> List[Path]:
    target_suffixes = {".xlsx", ".xlsm", ".xls", ".docx", ".doc", ".pdf", ".ppt", ".pptx", ".vsd", ".vsdx"}
    found: List[Path] = []
    exclude_resolved: Set[Path] = {p.resolve() for p in (exclude_paths or set())}
    ignored_dir_names = {".venv", "venv", ".git", "__pycache__"}
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in ignored_dir_names]
        dir_path = Path(dirpath)
        resolved_dir = dir_path.resolve()
        if resolved_dir in exclude_resolved or any(parent in exclude_resolved for parent in resolved_dir.parents):
            dirnames[:] = []
            continue
        for filename in filenames:
            if filename.startswith("~$"):
                continue
            path = dir_path / filename
            # Skip checker outputs to avoid recursively re-checking generated reports.
            if path.suffix.lower() == ".xlsx" and path.stem.lower().startswith("review_results"):
                continue
            if path.suffix.lower() not in target_suffixes:
                continue
            resolved = path.resolve()
            if resolved in exclude_resolved:
                continue
            if any(parent in exclude_resolved for parent in resolved.parents):
                continue
            found.append(path)
    return sorted(set(found))


def find_other_files(root: Path, exclude_paths: Optional[Set[Path]] = None) -> List[Path]:
    # チェック対象拡張子（Excel, Word, PDF, PPT, Visio）のみを厳密に除外
    target_suffixes = {".xlsx", ".xlsm", ".xls", ".docx", ".doc", ".pdf", ".ppt", ".pptx", ".vsd", ".vsdx"}
    found: List[Path] = []
    exclude_resolved: Set[Path] = {p.resolve() for p in (exclude_paths or set())}
    ignored_dir_names = {".venv", "venv", ".git", "__pycache__"}
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in ignored_dir_names]
        dir_path = Path(dirpath)
        resolved_dir = dir_path.resolve()
        if resolved_dir in exclude_resolved or any(parent in exclude_resolved for parent in resolved_dir.parents):
            dirnames[:] = []
            continue
        for filename in filenames:
            if filename.startswith("~$"):
                continue
            path = dir_path / filename
            if path.suffix.lower() == ".xlsx" and path.stem.lower().startswith("review_results"):
                continue
            resolved = path.resolve()
            if resolved in exclude_resolved:
                continue
            if any(parent in exclude_resolved for parent in resolved.parents):
                continue
            # チェック対象拡張子以外はother_filesに必ず含める
            if path.suffix.lower() not in target_suffixes:
                found.append(path)
            # それ以外（対象拡張子）はスキップ
    return sorted(set(found))


def extract_file_text_and_meta(file_path: Path) -> Tuple[str, Set[str], Set[str]]:
    suffix = file_path.suffix.lower()
    texts: List[str] = []
    annex_nums: Set[str] = set()
    wbs_nums: Set[str] = set()

    try:
        if suffix in {".xlsx", ".xlsm"} and load_workbook is not None:
            wb = None
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore", message=r"(?i)cannot parse header or footer so it will be ignored", category=UserWarning)
                wb = load_workbook(file_path, data_only=False)
            try:
                for text, _ in extract_excel_text_pages(wb):
                    texts.append(text)
            finally:
                if wb is not None:
                    try:
                        wb.close()
                    except Exception:
                        pass
        elif suffix == ".docx" and Document is not None:
            doc = Document(file_path)
            for p in doc.paragraphs:
                if p.text and p.text.strip():
                    texts.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if p.text and p.text.strip():
                                texts.append(p.text)
        elif suffix == ".pdf" and PdfReader is not None:
            reader = PdfReader(str(file_path))
            for text, _ in extract_pdf_text_pages(reader):
                texts.append(text)
        elif suffix in {".vsd", ".vsdx"}:
            diagram, err = load_visio_diagram_via_aspose(file_path)
            if diagram is not None:
                for text, _ in extract_visio_text_pages(diagram):
                    texts.append(text)
    except Exception:
        pass

    merged = "\n".join(texts)
    annex_nums |= parse_annex_numbers(file_path.stem)
    annex_nums |= parse_annex_numbers(merged)
    wbs_nums |= parse_wbs_numbers(merged)
    return merged, annex_nums, wbs_nums





def main(
    argv: Optional[List[str]] = None,
    progress_callback: Optional[Callable[[Dict[str, object]], None]] = None,
    cancel_requested: Optional[Callable[[], bool]] = None,
) -> None:
    parser = argparse.ArgumentParser(description="指定フォルダ配下のPDF/Word/Excelを共通+種別ルールでチェックし、結果xlsxを出力します。")
    parser.add_argument("root_folder", help="チェック対象のルートフォルダ")
    parser.add_argument("-o", "-0", "--output", default="review_results.xlsx", help="出力xlsxパス (-o と -0 の両方を使用可能)")
    parser.add_argument("--action-config-xlsx", default=None, help="対応推奨文言の上書き設定Excel。未指定時は既存の出力xlsx内設定シートを優先利用")
    parser.add_argument("--cover-keyword", default=None, help="表紙判定用キーワード")
    parser.add_argument("--visual-assets-dir", default="review_visual_assets", help="PDF/PNG生成物の保存先")
    parser.add_argument("--preview-pages-per-row", type=int, default=6, help="image_previewシートで1行に並べるページ数")
    parser.add_argument("--no-visual", action="store_true", help="画像生成（PDF/PNG・image_preview）を無効化する")
    args = parser.parse_args(argv)

    root = Path(args.root_folder).resolve()
    out_xlsx = Path(args.output).resolve()
    if not root.exists() or not root.is_dir():
        raise SystemExit(f"対象フォルダが存在しません: {root}")

    config_xlsx: Optional[Path] = Path(args.action_config_xlsx).resolve() if args.action_config_xlsx else (out_xlsx if out_xlsx.exists() else None)
    action_overrides, action_settings, action_settings_warning = load_suggested_action_settings(config_xlsx)

    assets_root = Path(args.visual_assets_dir).resolve()
    # パフォーマンス改善: ファイルリストを事前に全て取得し、進捗・残数を都度表示
    target_files = list(find_target_files(root, exclude_paths={out_xlsx, assets_root}))
    other_files = list(find_other_files(root, exclude_paths={out_xlsx, assets_root}))
    # チェック対象外ファイルはresultsから除外
    target_files = [f for f in target_files if f.suffix.lower() in {".xlsx", ".xlsm", ".xls", ".docx", ".doc", ".pdf", ".ppt", ".pptx", ".vsd", ".vsdx"}]
    # --- ファイルプロパティ情報取得関数 ---
    def get_file_property_detail(file_path: Path) -> str:
        """
        ファイルのプロパティ情報（タイトル、件名、タグ、分類、作成者、前回保存者、改訂番号、バージョン番号）を取得し、
        "タイトル=... / 件名=... / ..." の形式で返す。取得できない場合は空欄。
        """
        suffix = file_path.suffix.lower()
        props = {
            "タイトル": "",
            "件名": "",
            "タグ": "",
            "分類": "",
            "作成者": "",
            "前回保存者": "",
            "改訂番号": "",
            "バージョン番号": ""
        }
        try:
            if suffix == ".docx" and Document is not None:
                doc = Document(file_path)
                core = doc.core_properties
                props["タイトル"] = getattr(core, "title", "") or ""
                props["件名"] = getattr(core, "subject", "") or ""
                props["タグ"] = getattr(core, "keywords", "") or ""
                props["分類"] = getattr(core, "category", "") or ""
                props["作成者"] = getattr(core, "author", "") or ""
                props["前回保存者"] = getattr(core, "last_modified_by", "") or ""
                props["改訂番号"] = getattr(core, "revision", "") or ""
                props["バージョン番号"] = getattr(core, "version", "") or ""
            elif suffix in {".xlsx", ".xlsm", ".xls"} and load_workbook is not None:
                wb = load_workbook(file_path, read_only=True, data_only=True)
                core = getattr(wb, "properties", None)
                if core:
                    props["タイトル"] = getattr(core, "title", "") or ""
                    props["件名"] = getattr(core, "subject", "") or ""
                    props["タグ"] = getattr(core, "keywords", "") or ""
                    props["分類"] = getattr(core, "category", "") or ""
                    props["作成者"] = getattr(core, "creator", "") or ""
                    props["前回保存者"] = getattr(core, "lastModifiedBy", "") or ""
                    props["改訂番号"] = getattr(core, "revision", "") or ""
                    props["バージョン番号"] = getattr(core, "version", "") or ""
                wb.close()
            # 他形式（pptx, pdf, vsd等）は必要に応じて拡張
        except Exception:
            pass
        return " / ".join([f"{k}={v}" for k, v in props.items()])
    results: List[CheckResult] = []
    visual_pages: List[VisualPage] = []
    visual_enabled = not args.no_visual
    start_ts = time.perf_counter()
    total_files = len(target_files)
    completed_files = 0
    failed_files = 0
    processing_files = 0
    print(f"対象ファイル総数: {total_files}")

    if action_settings_warning:
        print(f"[WARN] {action_settings_warning}")
        add_result(results, root, "System", "S5", "対応推奨設定読込", "WARN", action_settings_warning, "設定Excelの suggested_action_settings シートを確認してください。")

    def report_progress(
        phase: str,
        current_file: str = "",
        cancelled: bool = False,
        file_index: Optional[int] = None,
        file_elapsed_sec: Optional[float] = None,
    ) -> None:
        if progress_callback is None:
            return
        # ファイル単位でOK/NG/ERROR/総数をカウントし、エラー内容も渡す
        file_status = {}
        file_errors = {}
        for r in results:
            fp = r.file_path
            if fp not in file_status:
                file_status[fp] = []
                file_errors[fp] = []
            file_status[fp].append(r.status)
            if r.status == "ERROR" and r.detail:
                file_errors[fp].append(r.detail)
        ok_count = 0
        ng_count = 0
        error_message = None
        for fp, statuses in file_status.items():
            if any(s == "FAIL" for s in statuses):
                ng_count += 1
            elif any(s == "ERROR" for s in statuses):
                ng_count += 1
                if not error_message and file_errors[fp]:
                    error_message = file_errors[fp][0]
            elif all(s == "PASS" for s in statuses):
                ok_count += 1
            # それ以外（MANUAL/WARN/N/Aのみ）はカウントしない
        all_count = total_files  # 総数は常に対象ファイル数
        payload = {
            "phase": phase,
            "completed": completed_files,
            "failed": failed_files,
            "processing": processing_files,
            "total": total_files,
            "current_file": current_file,
            "cancelled": cancelled,
            "ok_count": ok_count,
            "ng_count": ng_count,
            "all_count": all_count,
        }
        if file_index is not None:
            payload["file_index"] = file_index
        if file_elapsed_sec is not None:
            payload["file_elapsed_sec"] = file_elapsed_sec
        if current_file:
            current_statuses = [r.status for r in results if r.file_path == current_file]
            if any(status in {"FAIL", "ERROR"} for status in current_statuses):
                payload["file_result"] = "ng"
            elif current_statuses and all(status == "PASS" for status in current_statuses):
                payload["file_result"] = "ok"
            elif current_statuses:
                payload["file_result"] = "manual"
        if error_message:
            payload["error_message"] = error_message
        progress_callback(payload)

    report_progress("start")
    print(f"処理開始: {time.strftime('%Y-%m-%d %H:%M:%S')}")

    if not target_files:
        add_result(results, root, "System", "S1", "対象ファイル検出", "WARN", "対象ファイル（Excel/Word/PDF）が見つかりません。", "フォルダ構成を確認。")


    else:
        cancelled = False
        for idx, file_path in enumerate(target_files, start=1):
            if cancel_requested is not None and cancel_requested():
                cancelled = True
                add_result(results, root, "System", "S4", "ユーザー中止", "WARN", "ユーザー操作により処理を中止しました。", "必要に応じて再実行してください。")
                break

            file_start = time.perf_counter()
            processing_files = 1
            report_progress("processing", current_file=str(file_path), file_index=idx)
            print(f"[{idx}/{total_files}] 開始: {display_file_path_for_log(file_path)}")
            try:
                try:
                    check_file(file_path, results, args.cover_keyword)
                    page_count = None
                    if visual_enabled:
                        page_count = run_visual_pipeline(file_path, results, visual_pages, assets_root)
                    append_max_page_detail(results, file_path, page_count)
                except Exception as exc:
                    failed_files += 1
                    import traceback
                    tb = traceback.format_exc()
                    add_result(results, file_path, "System", "S3", "ファイル処理例外", "ERROR", f"処理中に例外が発生しました: {exc}\n{tb}", "対象ファイルを個別確認し、必要に応じて再実行してください。")
            finally:
                elapsed = time.perf_counter() - file_start
                completed_files += 1
                processing_files = 0
                report_progress("processed", current_file=str(file_path), file_index=idx, file_elapsed_sec=elapsed)
                print(f"[{idx}/{total_files}] 完了: {display_file_path_for_log(file_path)} ({elapsed:.2f}秒)")


    # 必ず統合ファイルのみ出力。分割ファイルは出力しない。
    apply_suggested_action_overrides(results, action_overrides)
    output_action_settings = build_suggested_action_settings(results, action_settings)
    err = write_visual_report_xlsx(
        results,
        visual_pages,
        out_xlsx,
        other_files=other_files,
        preview_pages_per_row=1,
        suggested_action_settings=output_action_settings,
    )
    if err:
        print(f"[ERROR] 統合xlsx出力失敗: {err}")

    total = len(results)
    fail = sum(1 for r in results if r.status == "FAIL")
    warn = sum(1 for r in results if r.status == "WARN")
    manual = sum(1 for r in results if r.status == "MANUAL")
    error = sum(1 for r in results if r.status == "ERROR")
    na = sum(1 for r in results if r.status == "N/A")

    print(f"対象ファイル数: {len(target_files)}")
    print(f"チェック結果数: {total} (FAIL={fail}, WARN={warn}, MANUAL={manual}, ERROR={error}, N/A={na})")
    print(f"xlsx出力: {out_xlsx}")
    print(f"画像生成: {'有効' if visual_enabled else '無効'}")
    print(f"処理終了: {time.strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"全体処理時間: {time.perf_counter() - start_ts:.2f}秒")
    report_progress("done", cancelled=(cancel_requested is not None and cancel_requested()))


if __name__ == "__main__":
    main()
